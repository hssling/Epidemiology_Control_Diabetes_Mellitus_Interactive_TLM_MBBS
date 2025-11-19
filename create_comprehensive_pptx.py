"""
Comprehensive Diabetes Mellitus Teaching Learning Material - PowerPoint Presentation Generator
Creates a detailed PPTX file covering all aspects of diabetes mellitus for medical education

To use this script:
1. Install required libraries: pip install python-pptx pandas matplotlib numpy plotly seaborn
2. Run: python create_comprehensive_pptx.py
3. The presentation will be saved as 'Diabetes_Comprehensive_TLM_Presentation.pptx'
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData
    import pandas as pd
    import matplotlib.pyplot as plt
    import numpy as np
    import seaborn as sns
    import os
    import io
except ImportError as e:
    print(f"Missing required libraries. Please install: pip install python-pptx pandas matplotlib numpy seaborn plotly")
    print(f"Error: {e}")
    exit(1)

# Set matplotlib style
plt.style.use('seaborn-v0_8-whitegrid')
sns.set_palette("husl")

def create_title_slide(prs):
    """Create comprehensive title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Diabetes Mellitus: Comprehensive Teaching Learning Material"
    subtitle.text = "Epidemiology, Pathophysiology, Management & Prevention\nMBBS Medical Education | Latest Research 2024-2025"

    # Style the title
    title_tf = title.text_frame
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    title_tf.paragraphs[0].font.bold = True

    subtitle_tf = subtitle.text_frame
    subtitle_tf.paragraphs[0].font.size = Pt(20)
    subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(52, 73, 94)

def create_overview_slide(prs):
    """Create comprehensive presentation overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title_shape = shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Comprehensive Diabetes Overview"
    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].font.size = Pt(36)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = ""

    sections = [
        "• DEFINITION & CLASSIFICATION - ADA/WHO Criteria 2024",
        "• GLOBAL EPIDEMIOLOGY - 537M cases, 11.3% prevalence",
        "• PATHOGENESIS - β-cell dysfunction, insulin resistance",
        "• CLINICAL FEATURES - Symptoms, complications spectrum",
        "• DIAGNOSIS - FPG, OGTT, HbA1c thresholds",
        "• MANAGEMENT - Lifestyle, pharmacological, technology",
        "• CONTROL STRATEGIES - Glycemic targets, monitoring",
        "• PREVENTION - DPP model, lifestyle intervention",
        "• NATIONAL PROGRAMME - NPCDCS implementation",
        "• FUTURE DIRECTIONS - Precision medicine, AI, stem cells"
    ]

    for section in sections:
        p = tf.add_paragraph()
        p.text = section
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(52, 73, 94)
        p.level = 0

    # Add learning objectives box
    obj_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    obj_tf = obj_box.text_frame
    obj_tf.text = "LEARNING OBJECTIVES: Define, diagnose, manage diabetes comprehensively; Apply evidence-based prevention strategies; Understand national programmes and future innovations"
    obj_tf.paragraphs[0].font.size = Pt(12)
    obj_tf.paragraphs[0].font.bold = True
    obj_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

def create_definitions_slide(prs):
    """Definition and classification slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Diabetes Mellitus: Definition & Classification (ADA 2024)"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # WHO Definition
    who_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(2))
    who_tf = who_box.text_frame
    who_tf.text = "WORLD HEALTH ORGANIZATION DEFINITION"
    who_tf.paragraphs[0].font.bold = True
    who_tf.paragraphs[0].font.size = Pt(16)
    who_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    who_points = [
        "• Chronic hyperglycemia due to defects in:",
        "  - Insulin secretion (β-cell dysfunction)",
        "  - Insulin action (insulin resistance)",
        "  - Both mechanisms",
        "• Characterized by disturbances in:",
        "  - Carbohydrate, fat, and protein metabolism"
    ]

    for point in who_points:
        p = who_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(12)

    # ADA Classification
    ada_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(2))
    ada_tf = ada_box.text_frame
    ada_tf.text = "ADA CLASSIFICATION SYSTEM"
    ada_tf.paragraphs[0].font.bold = True
    ada_tf.paragraphs[0].font.size = Pt(16)
    ada_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    ada_types = [
        "• Type 1 DM: Autoimmune β-cell destruction",
        "• Type 2 DM: Insulin resistance + relative insulin deficiency",
        "• Gestational DM: Glucose intolerance in pregnancy",
        "• Other specific types: Monogenic, drug-induced, etc.",
        "• Prediabetes: IFG, IGT, or both"
    ]

    for type_ in ada_types:
        p = ada_tf.add_paragraph()
        p.text = type_
        p.font.size = Pt(12)

    # Historical perspective
    hist_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.5))
    hist_tf = hist_box.text_frame
    hist_tf.text = "HISTORICAL EVOLUTION: Ancient descriptions → 1921 Insulin discovery → 1950s Oral agents → 1980s HbA1c → 2020s CGM/Precision medicine"
    hist_tf.paragraphs[0].font.size = Pt(14)
    hist_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    hist_tf.paragraphs[0].font.bold = True

def create_epidemiology_global_slide(prs):
    """Global epidemiology slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Global Diabetes Epidemiology (IDF 2021)"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Key statistics
    stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(4))
    stats_tf = stats_box.text_frame
    stats_tf.text = "GLOBAL BURDEN STATISTICS"
    stats_tf.paragraphs[0].font.bold = True
    stats_tf.paragraphs[0].font.size = Pt(16)
    stats_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

    stats = [
        "• 537 MILLION adults with diabetes (2021)",
        "• 1 in 10 adults affected worldwide",
        "• 541 million in prediabetes stage",
        "• 75% live in low- and middle-income countries",
        "• 6.7 million deaths annually",
        "• $966 billion healthcare expenditure",
        "• Projected: 783 million by 2045"
    ]

    for stat in stats:
        p = stats_tf.add_paragraph()
        p.text = stat
        p.font.size = Pt(12)
        p.font.bold = True

    # Regional variations
    regions_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(4))
    regions_tf = regions_box.text_frame
    regions_tf.text = "REGIONAL PREVALENCE (%)"
    regions_tf.paragraphs[0].font.bold = True
    regions_tf.paragraphs[0].font.size = Pt(16)
    regions_tf.paragraphs[0].font.color.rgb = RGBColor(155, 89, 182)

    regions = [
        "• Middle East & North Africa: 14.0%",
        "• North America: 11.7%",
        "• South & Central America: 9.9%",
        "• South Asia: 8.8%",
        "• East Asia & Pacific: 8.4%",
        "• Western Europe: 5.4%",
        "• Sub-Saharan Africa: 3.3%"
    ]

    for region in regions:
        p = regions_tf.add_paragraph()
        p.text = region
        p.font.size = Pt(12)

    # Trend projection
    trend_box = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(6), Inches(0.8))
    trend_tf = trend_box.text_frame
    trend_tf.text = "TREND: 3-fold increase since 1980; Slowing growth in high-income countries; Rapid rise in LMICs"
    trend_tf.paragraphs[0].font.size = Pt(12)
    trend_tf.paragraphs[0].font.bold = True
    trend_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_epidemiology_india_slide(prs):
    """India-specific epidemiology slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Diabetes Epidemiology in India"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # India statistics
    india_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(3))
    india_tf = india_box.text_frame
    india_tf.text = "INDIA DIABETES BURDEN (2021)"
    india_tf.paragraphs[0].font.bold = True
    india_tf.paragraphs[0].font.size = Pt(16)
    india_tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)

    india_stats = [
        "• 101.2 MILLION people with diabetes",
        "• 11.4% prevalence (20-79 years)",
        "• Second highest globally after China",
        "• 136 million with prediabetes",
        "• Annual cost: INR 2.5-3 lakh crores",
        "• Projected: 140 million by 2030"
    ]

    for stat in india_stats:
        p = india_tf.add_paragraph()
        p.text = stat
        p.font.size = Pt(12)
        p.font.bold = True

    # Regional variations
    regional_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(3))
    regional_tf = regional_box.text_frame
    regional_tf.text = "REGIONAL VARIATIONS IN INDIA"
    regional_tf.paragraphs[0].font.bold = True
    regional_tf.paragraphs[0].font.size = Pt(16)
    regional_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    regional_stats = [
        "• Kerala, Punjab, Tamil Nadu: >15%",
        "• Andhra Pradesh, Karnataka: 13-15%",
        "• Maharashtra, Gujarat: 10-13%",
        "• Rural areas: Lower prevalence",
        "• Urban-rural ratio: 2:1"
    ]

    for stat in regional_stats:
        p = regional_tf.add_paragraph()
        p.text = stat
        p.font.size = Pt(12)

    # Contributing factors
    factors_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
    factors_tf = factors_box.text_frame
    factors_tf.text = "KEY CONTRIBUTING FACTORS: Urbanization, lifestyle changes, genetic predisposition, low birth weight, dietary transition, physical inactivity"
    factors_tf.paragraphs[0].font.size = Pt(12)
    factors_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    factors_tf.paragraphs[0].font.bold = True

def create_epidemiology_visual_slide(prs):
    """Epidemiology with visual chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Diabetes Trends: Global vs India"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add epidemiology chart
    try:
        chart_path = "TLM_Diabetes_Mellitus/visualizations/epidemiology_chart.png"
        slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    except:
        # Fallback statistics
        fallback_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
        fallback_tf = fallback_box.text_frame
        fallback_tf.text = "Epidemiology chart not available. Key trends:\n\nGlobal: 537M cases (2021) → 783M (2045)\nIndia: 101M cases (2021) → 160M (2045)\n\nRegional variations: Middle East highest (14%), Africa lowest (3.3%)\nUrban-rural ratio in LMICs: 2-3:1"
        for para in fallback_tf.paragraphs:
            para.font.size = Pt(16)

def create_pathophysiology_slide(prs):
    """Pathophysiology overview slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Pathophysiology of Diabetes Mellitus"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Normal glucose metabolism
    normal_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.6), Inches(2.5))
    normal_tf = normal_box.text_frame
    normal_tf.text = "NORMAL GLUCOSE METABOLISM"
    normal_tf.paragraphs[0].font.bold = True
    normal_tf.paragraphs[0].font.size = Pt(14)
    normal_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    normal_points = [
        "• Insulin secretion (β-cells)",
        "• Insulin action (target tissues)",
        "• Glucose uptake (GLUT4 translocation)",
        "• Glycogen synthesis",
        "• Gluconeogenesis suppression"
    ]

    for point in normal_points:
        p = normal_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

    # Type 1 DM
    t1_box = slide.shapes.add_textbox(Inches(3.3), Inches(1.2), Inches(2.6), Inches(2.5))
    t1_tf = t1_box.text_frame
    t1_tf.text = "TYPE 1 DM PATHOGENESIS"
    t1_tf.paragraphs[0].font.bold = True
    t1_tf.paragraphs[0].font.size = Pt(14)
    t1_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

    t1_points = [
        "• Autoimmune destruction of β-cells",
        "• Genetic susceptibility (HLA-DR3/4)",
        "• Environmental triggers (virus?)",
        "• Insulitis → absolute insulin deficiency",
        "• Ketoacidosis risk high"
    ]

    for point in t1_points:
        p = t1_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

    # Type 2 DM
    t2_box = slide.shapes.add_textbox(Inches(6.1), Inches(1.2), Inches(2.6), Inches(2.5))
    t2_tf = t2_box.text_frame
    t2_tf.text = "TYPE 2 DM PATHOGENESIS"
    t2_tf.paragraphs[0].font.bold = True
    t2_tf.paragraphs[0].font.size = Pt(14)
    t2_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    t2_points = [
        "• Insulin resistance (primary defect)",
        "• β-cell dysfunction (secondary)",
        "• Increased gluconeogenesis",
        "• Decreased insulin secretion",
        "• Chronic inflammation"
    ]

    for point in t2_points:
        p = t2_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

def create_pathophysiology_visual_slide(prs):
    """Pathophysiology visual diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Molecular Mechanisms in Diabetes Pathogenesis"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add pathophysiology diagram
    try:
        diagram_path = "TLM_Diabetes_Mellitus/visualizations/pathophysiology_diagram.png"
        slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    except:
        # Fallback text describing mechanisms
        mech_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
        mech_tf = mech_box.text_frame
        mech_tf.text = "KEY PATHOGENIC MECHANISMS:\n\n1. INSULIN RESISTANCE:\n   • Impaired IRS-1/PI3K signaling\n   • Increased FFA flux (adipose tissue)\n   • Inflammation (TNF-α, IL-6)\n   • Mitochondrial dysfunction\n\n2. β-CELL DYSFUNCTION:\n   • Amyloid deposition\n   • Oxidative stress\n   • Glucotoxicity & lipotoxicity\n   • Impaired insulin processing\n\n3. LIVER: Increased gluconeogenesis (PEPCK, G6Pase upregulation)\n\n4. MUSCLE: Reduced GLUT4 translocation, glycogen synthesis\n\n5. ADIPOSE: Adipokine imbalance (↓adiponectin, ↑leptin resistance)"
        mech_tf.paragraphs[0].font.size = Pt(12)
        mech_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_types_diabetes_slide(prs):
    """Types of diabetes comprehensive slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Classification of Diabetes Mellitus"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Type 1 DM
    t1_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.2), Inches(4))
    t1_tf = t1_box.text_frame
    t1_tf.text = "TYPE 1 DM\n(5-10%)"
    t1_tf.paragraphs[0].font.bold = True
    t1_tf.paragraphs[0].font.size = Pt(14)
    t1_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

    t1_features = [
        "• Age: Usually <30 years",
        "• Onset: Acute/rapid",
        "• Autoantibodies: GAD, IA2",
        "• C-peptide: Low/absent",
        "• HLA association: DR3/4",
        "• Complications: Ketoacidosis common"
    ]

    for feature in t1_features:
        p = t1_tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(10)

    # Type 2 DM
    t2_box = slide.shapes.add_textbox(Inches(2.9), Inches(1.2), Inches(2.2), Inches(4))
    t2_tf = t2_box.text_frame
    t2_tf.text = "TYPE 2 DM\n(90-95%)"
    t2_tf.paragraphs[0].font.bold = True
    t2_tf.paragraphs[0].font.size = Pt(14)
    t2_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    t2_features = [
        "• Age: Usually >40 years",
        "• Onset: Insidious/gradual",
        "• Insulin resistance: Primary",
        "• C-peptide: Initially normal",
        "• Obesity: Common",
        "• Family history: Strong"
    ]

    for feature in t2_features:
        p = t2_tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(10)

    # Other types
    other_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.2), Inches(2.2), Inches(4))
    other_tf = other_box.text_frame
    other_tf.text = "OTHER SPECIFIC TYPES"
    other_tf.paragraphs[0].font.bold = True
    other_tf.paragraphs[0].font.size = Pt(14)
    other_tf.paragraphs[0].font.color.rgb = RGBColor(155, 89, 182)

    other_types = [
        "• MODY: Maturity onset",
        "• LADA: Latent autoimmune",
        "• Neonatal DM",
        "• Drug-induced (steroids)",
        "• Pancreatic diseases",
        "• Genetic syndromes"
    ]

    for type_ in other_types:
        p = other_tf.add_paragraph()
        p.text = type_
        p.font.size = Pt(10)

    # Gestational and Prediabetes
    gdm_box = slide.shapes.add_textbox(Inches(7.7), Inches(1.2), Inches(2.2), Inches(4))
    gdm_tf = gdm_box.text_frame
    gdm_tf.text = "GESTATIONAL & PREDIABETES"
    gdm_tf.paragraphs[0].font.bold = True
    gdm_tf.paragraphs[0].font.size = Pt(14)
    gdm_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    gdm_types = [
        "• GDM: Pregnancy onset",
        "• Prediabetes: IFG/IGT",
        "• Risk of T2DM: High",
        "• Screening: Universal",
        "• Management: Lifestyle ± metformin",
        "• Postpartum monitoring"
    ]

    for type_ in gdm_types:
        p = gdm_tf.add_paragraph()
        p.text = type_
        p.font.size = Pt(10)

def create_risk_factors_slide(prs):
    """Risk factors comprehensive slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Risk Factors for Diabetes Mellitus"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Non-modifiable factors
    nonmod_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3), Inches(2.5))
    nonmod_tf = nonmod_box.text_frame
    nonmod_tf.text = "NON-MODIFIABLE RISK FACTORS"
    nonmod_tf.paragraphs[0].font.bold = True
    nonmod_tf.paragraphs[0].font.size = Pt(14)
    nonmod_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

    nonmod_factors = [
        "• Age >45 years",
        "• Family history (first-degree)",
        "• Genetics (T2DM: polygenic)",
        "• Ethnicity (South Asians high risk)",
        "• Previous GDM",
        "• Low birth weight"
    ]

    for factor in nonmod_factors:
        p = nonmod_tf.add_paragraph()
        p.text = factor
        p.font.size = Pt(11)

    # Modifiable lifestyle factors
    lifestyle_box = slide.shapes.add_textbox(Inches(3.8), Inches(1.2), Inches(3), Inches(2.5))
    lifestyle_tf = lifestyle_box.text_frame
    lifestyle_tf.text = "MODIFIABLE LIFESTYLE FACTORS"
    lifestyle_tf.paragraphs[0].font.bold = True
    lifestyle_tf.paragraphs[0].font.size = Pt(14)
    lifestyle_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    lifestyle_factors = [
        "• Obesity (BMI >25 kg/m²)",
        "• Central obesity (waist ≥90cm M, ≥80cm F)",
        "• Physical inactivity (<150 min/week)",
        "• Unhealthy diet (high refined carbs)",
        "• Smoking and alcohol excess",
        "• Sleep deprivation (<7 hours)"
    ]

    for factor in lifestyle_factors:
        p = lifestyle_tf.add_paragraph()
        p.text = factor
        p.font.size = Pt(11)

    # Medical conditions
    medical_box = slide.shapes.add_textbox(Inches(7), Inches(1.2), Inches(2.75), Inches(2.5))
    medical_tf = medical_box.text_frame
    medical_tf.text = "MEDICAL CONDITIONS"
    medical_tf.paragraphs[0].font.bold = True
    medical_tf.paragraphs[0].font.size = Pt(14)
    medical_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    medical_conditions = [
        "• Polycystic ovary syndrome",
        "• Metabolic syndrome",
        "• Hypertension",
        "• Dyslipidemia",
        "• NAFLD",
        "• Psychiatric disorders"
    ]

    for condition in medical_conditions:
        p = medical_tf.add_paragraph()
        p.text = condition
        p.font.size = Pt(11)

    # Indian risk assessment
    india_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    india_tf = india_box.text_frame
    india_tf.text = "INDIAN DIABETES RISK SCORE (IDRS): Age + BMI + Waist + FH + Activity. Score ≥60 = High risk. Validated for Asian Indians."
    india_tf.paragraphs[0].font.size = Pt(12)
    india_tf.paragraphs[0].font.bold = True
    india_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_risk_factors_visual_slide(prs):
    """Risk factors visual diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Risk Factor Interactions & Prevention Opportunities"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add risk factors diagram
    try:
        diagram_path = "TLM_Diabetes_Mellitus/visualizations/risk_factor_diagram.png"
        slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    except:
        # Fallback text
        fallback_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
        fallback_tf = fallback_box.text_frame
        fallback_tf.text = "RISK FACTOR MODEL:\n\n  Genetics + Environment → Insulin Resistance\n                ↓\nβ-cell Compensation → Failure → Hyperglycemia\n                ↓\nT2DM Development → Micro/Macrovascular Complications\n\nPREVENTION OPPORTUNITIES:\n• Primary: Lifestyle intervention at risk stage\n• Secondary: Early detection and intervention\n• Tertiary: Optimal control to prevent complications"
        for para in fallback_tf.paragraphs:
            para.font.size = Pt(16)

def create_clinical_features_slide(prs):
    """Clinical features and signs"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Clinical Features of Diabetes Mellitus"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Classic symptoms (3P's)
    symptoms_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3), Inches(2.5))
    symptoms_tf = symptoms_box.text_frame
    symptoms_tf.text = "CLASSIC SYMPTOMS (3P's)"
    symptoms_tf.paragraphs[0].font.bold = True
    symptoms_tf.paragraphs[0].font.size = Pt(14)
    symptoms_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

    symptoms = [
        "• Polyuria (increased urination)",
        "• Polydipsia (increased thirst)",
        "• Polyphagia (increased hunger)",
        "• Unexplained weight loss",
        "• Fatigue and weakness",
        "• Blurred vision"
    ]

    for symptom in symptoms:
        p = symptoms_tf.add_paragraph()
        p.text = symptom
        p.font.size = Pt(11)

    # Asymptomatic presentation
    asymp_box = slide.shapes.add_textbox(Inches(3.8), Inches(1.2), Inches(3), Inches(2.5))
    asymp_tf = asymp_box.text_frame
    asymp_tf.text = "ASYMPTOMATIC PRESENTATION"
    asymp_tf.paragraphs[0].font.bold = True
    asymp_tf.paragraphs[0].font.size = Pt(14)
    asymp_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    asymp_features = [
        "• Routine screening detection",
        "• Incidental hyperglycemia",
        "• Complications presentation",
        "• Associated conditions",
        "• Type 2 DM common pattern",
        "• Delayed diagnosis risks"
    ]

    for feature in asymp_features:
        p = asymp_tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(11)

    # Physical examination findings
    exam_box = slide.shapes.add_textbox(Inches(7), Inches(1.2), Inches(2.75), Inches(2.5))
    exam_tf = exam_box.text_frame
    exam_tf.text = "PHYSICAL EXAMINATION"
    exam_tf.paragraphs[0].font.bold = True
    exam_tf.paragraphs[0].font.size = Pt(14)
    exam_tf.paragraphs[0].font.color.rgb = RGBColor(155, 89, 182)

    exam_findings = [
        "• BMI and waist circumference",
        "• Blood pressure measurement",
        "• Signs of dehydration",
        "• Insulin resistance markers",
        "• Foot examination",
        "• Oral health assessment"
    ]

    for finding in exam_findings:
        p = exam_tf.add_paragraph()
        p.text = finding
        p.font.size = Pt(11)

    # Acute vs Chronic presentations
    pres_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    pres_tf = pres_box.text_frame
    pres_tf.text = "DIFFERENTIAL PRESENTATION: T1DM often acute with DKA risk; T2DM usually insidious with complications; GDM detected during pregnancy screening; MODY may present with mild symptoms"
    pres_tf.paragraphs[0].font.size = Pt(12)
    pres_tf.paragraphs[0].font.bold = True
    pres_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_diagnosis_criteria_slide(prs):
    """Diagnosis criteria and methods"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Diagnosis of Diabetes Mellitus (ADA 2024)"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Diagnostic criteria table
    criteria_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(4))
    criteria_tf = criteria_box.text_frame
    criteria_tf.text = "DIAGNOSTIC THRESHOLDS"
    criteria_tf.paragraphs[0].font.bold = True
    criteria_tf.paragraphs[0].font.size = Pt(16)
    criteria_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    criteria = [
        "• HbA1c ≥ 6.5% (NGSP certified)",
        "• FPG ≥ 126 mg/dL (8.0 mmol/L)",
        "• 2h-PG ≥ 200 mg/dL (11.1 mmol/L)",
        "• Random PG ≥ 200 mg/dL + symptoms",
        "• OGTT: 75g glucose load",
        "• Repeat testing for confirmation"
    ]

    for criterion in criteria:
        p = criteria_tf.add_paragraph()
        p.text = criterion
        p.font.size = Pt(12)

    # Prediabetes criteria
    prediabetes_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4), Inches(4))
    prediabetes_tf = prediabetes_box.text_frame
    prediabetes_tf.text = "PREDIABETES CRITERIA"
    prediabetes_tf.paragraphs[0].font.bold = True
    prediabetes_tf.paragraphs[0].font.size = Pt(16)
    prediabetes_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    prediabetes_criteria = [
        "• IFG: FPG 100-125 mg/dL (5.6-6.9 mmol/L)",
        "• IGT: 2h-PG 140-199 mg/dL (7.8-11.0 mmol/L)",
        "• HbA1c 5.7-6.4% (may be included)",
        "• High risk for T2DM progression",
        "• Requires lifestyle intervention",
        "• Annual monitoring recommended"
    ]

    for criterion in prediabetes_criteria:
        p = prediabetes_tf.add_paragraph()
        p.text = criterion
        p.font.size = Pt(12)

    # Special situations note
    special_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.8))
    special_tf = special_box.text_frame
    special_tf.text = "SPECIAL SITUATIONS: Pregnancy (IADPSG criteria), children (pediatric thresholds), symptomatic patients (random glucose)"
    special_tf.paragraphs[0].font.size = Pt(12)
    special_tf.paragraphs[0].font.bold = True
    special_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_treatment_overview_slide(prs):
    """Treatment and management overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Management of Diabetes Mellitus"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Non-pharmacological management
    nonpharm_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3.5), Inches(3))
    nonpharm_tf = nonpharm_box.text_frame
    nonpharm_tf.text = "NON-PHARMACOLOGICAL MANAGEMENT"
    nonpharm_tf.paragraphs[0].font.bold = True
    nonpharm_tf.paragraphs[0].font.size = Pt(14)
    nonpharm_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    nonpharm = [
        "• Diabetes Self-Management Education (DSME)",
        "• Medical Nutrition Therapy (MNT)",
        "• Physical Activity (150 min/week)",
        "• Weight Management (5-10% weight loss)",
        "• Smoking Cessation",
        "• Problem Solving Support"
    ]

    for item in nonpharm:
        p = nonpharm_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)

    # Pharmacological management
    pharm_box = slide.shapes.add_textbox(Inches(4.5), Inches(1.2), Inches(5), Inches(3))
    pharm_tf = pharm_box.text_frame
    pharm_tf.text = "PHARMACOLOGICAL MANAGEMENT"
    pharm_tf.paragraphs[0].font.bold = True
    pharm_tf.paragraphs[0].font.size = Pt(14)
    pharm_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    pharm_agents = [
        "• Metformin: First-line for T2DM",
        "• SGLT2i: CV and renal benefits (empagliflozin)",
        "• GLP-1 RA: Weight loss and CV protection (semaglutide)",
        "• DPP-4i: Low hypoglycemia risk (sitagliptin)",
        "• Insulin: Essential for T1DM, progressive T2DM",
        "• TZDs: Insulin sensitizers (pioglitazone)"
    ]

    for agent in pharm_agents:
        p = pharm_tf.add_paragraph()
        p.text = agent
        p.font.size = Pt(11)

    # Insulin regimens
    insulin_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2))
    insulin_tf = insulin_box.text_frame
    insulin_tf.text = "INSULIN THERAPY: Basal (long-acting) + Bolus (rapid-acting). Regimens: Basal-bolus, premixed, basal-plus. Individualized based on needs, lifestyle, and response."
    insulin_tf.paragraphs[0].font.size = Pt(12)
    insulin_tf.paragraphs[0].font.bold = True
    insulin_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_treatment_algorithm_slide(prs):
    """Treatment algorithm visual"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "ADA Treatment Algorithm for Type 2 Diabetes"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add treatment algorithm diagram
    try:
        diagram_path = "TLM_Diabetes_Mellitus/visualizations/treatment_algorithm.png"
        slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    except:
        # Fallback algorithm text
        algo_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
        algo_tf = algo_box.text_frame
        algo_tf.text = "ADA TREATMENT ALGORITHM FOR T2DM:\n\n1. Diagnosis → Lifestyle + Metformin\n\n2. Metformin monotherapy → Add second agent:\n   • CV risk high: SGLT2i or GLP-1 RA\n   • CV risk low + weight concern: GLP-1 RA\n   • No weight concern: DPP-4i or SGLT2i\n\n3. Triple therapy → Add insulin if needed\n\n4. Insulin → Intensive regimen (basal-bolus)\n\n5. Consider comorbidities: CKD (avoid metformin), HF (SGLT2i)"
        algo_tf.paragraphs[0].font.size = Pt(12)
        algo_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_control_targets_slide(prs):
    """Glycemic control targets (updated with latest research)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Glycemic Control Targets (ADA 2024)"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # General targets
    general_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(3))
    general_tf = general_box.text_frame
    general_tf.text = "GENERAL TARGETS"
    general_tf.paragraphs[0].font.bold = True
    general_tf.paragraphs[0].font.size = Pt(18)
    general_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    general_targets = [
        "• HbA1c: <7.0% (most adults)",
        "• Preprandial glucose: 80-130 mg/dL",
        "• Peak postprandial: <180 mg/dL",
        "• Bedtime glucose: 100-140 mg/dL",
        "• TIR (Time in Range): >70%"
    ]

    for target in general_targets:
        p = general_tf.add_paragraph()
        p.text = target
        p.font.size = Pt(16)

    # Individualized targets
    individual_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(3))
    individual_tf = individual_box.text_frame
    individual_tf.text = "INDIVIDUALIZED TARGETS"
    individual_tf.paragraphs[0].font.bold = True
    individual_tf.paragraphs[0].font.size = Pt(18)
    individual_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    individual_targets = [
        "• Young patients (<40): HbA1c <6.5-7.0%",
        "• Elderly (≥65): HbA1c 7.5-8.0%",
        "• Longstanding DM (>10y): More liberal",
        "• High CV risk: Stringent control",
        "• Frequent hypoglycemia: Relaxed targets"
    ]

    for target in individual_targets:
        p = individual_tf.add_paragraph()
        p.text = target
        p.font.size = Pt(16)

    # Latest evidence note
    evidence_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    evidence_tf = evidence_box.text_frame
    evidence_tf.text = "LATEST EVIDENCE (2024): TIR >70% and HbA1c <7% balance benefits vs. hypoglycemia. GRADE study: Intensive control (HbA1c <6.5%) prevents complications but increases hypoglycemia. Personalized targets based on CGM, comorbidities, and patient preferences."
    evidence_tf.paragraphs[0].font.size = Pt(14)
    evidence_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_monitoring_slide(prs):
    """Comprehensive monitoring strategies"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Diabetes Monitoring Strategies (2024)"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # SMBG recommendations
    smbg_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.8), Inches(2.5))
    smbg_tf = smbg_box.text_frame
    smbg_tf.text = "SELF-MONITORING OF BLOOD GLUCOSE"
    smbg_tf.paragraphs[0].font.bold = True
    smbg_tf.paragraphs[0].font.size = Pt(14)
    smbg_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    smbg_points = [
        "• T2DM oral agents: 1-2x daily",
        "• T2DM on insulin: 2-4x daily",
        "• T1DM: 4-6x daily (basal-bolus)",
        "• Pre/postprandial monitoring",
        "• Pattern recognition"
    ]

    for point in smbg_points:
        p = smbg_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

    # CGM recommendations
    cgm_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.2), Inches(2.8), Inches(2.5))
    cgm_tf = cgm_box.text_frame
    cgm_tf.text = "CONTINUOUS GLUCOSE MONITORING"
    cgm_tf.paragraphs[0].font.bold = True
    cgm_tf.paragraphs[0].font.size = Pt(14)
    cgm_tf.paragraphs[0].font.color.rgb = RGBColor(155, 89, 182)

    cgm_points = [
        "• TIR (>70%): Primary outcome",
        "• TBR (<4%): Minimize hypo",
        "• TAR (<25%): Manage hyper",
        "• AGP analysis",
        "• Real-time vs. intermittent"
    ]

    for point in cgm_points:
        p = cgm_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

    # Biomarkers for complications
    biomarkers_box = slide.shapes.add_textbox(Inches(6.3), Inches(1.2), Inches(3), Inches(2.5))
    biomarkers_tf = biomarkers_box.text_frame
    biomarkers_tf.text = "COMPREHENSIVE MONITORING"
    biomarkers_tf.paragraphs[0].font.bold = True
    biomarkers_tf.paragraphs[0].font.size = Pt(14)
    biomarkers_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    biomarkers = [
        "• HbA1c: Every 3-6 months",
        "• Lipid profile: Annual",
        "• Kidney function: eGFR, ACR",
        "• Eye examination: Annual",
        "• Foot assessment: Annual"
    ]

    for biomarker in biomarkers:
        p = biomarkers_tf.add_paragraph()
        p.text = biomarker
        p.font.size = Pt(11)

    # Technology advances
    tech_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    tech_tf = tech_box.text_frame
    tech_tf.text = "2024 ADVANCES: AI-powered CGM with predictive alerts, integrated digital platforms, remote monitoring capabilities, automated insulin delivery systems"
    tech_tf.paragraphs[0].font.size = Pt(12)
    tech_tf.paragraphs[0].font.bold = True
    tech_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_prevention_overview_slide(prs):
    """Prevention strategies overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Diabetes Prevention: Evidence-Based Strategies"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Prevention levels
    primary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.8), Inches(2.5))
    primary_tf = primary_box.text_frame
    primary_tf.text = "PRIMARY PREVENTION\n(Prevent Onset)"
    primary_tf.paragraphs[0].font.bold = True
    primary_tf.paragraphs[0].font.size = Pt(14)
    primary_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

    primary_points = [
        "• Target: High-risk prediabetics",
        "• DPP model: 58% reduction",
        "• Intensive lifestyle intervention",
        "• 5-7% weight loss",
        "• 150 min/week activity",
        "• Metformin optional (BMI ≥35)"
    ]

    for point in primary_points:
        p = primary_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(10)

    # Secondary prevention
    secondary_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.2), Inches(2.8), Inches(2.5))
    secondary_tf = secondary_box.text_frame
    secondary_tf.text = "SECONDARY PREVENTION\n(Delay Progression)"
    secondary_tf.paragraphs[0].font.bold = True
    secondary_tf.paragraphs[0].font.size = Pt(14)
    secondary_tf.paragraphs[0].font.color.rgb = RGBColor(243, 156, 18)

    secondary_points = [
        "• Prediabetes management",
        "• IFG: 100-125 mg/dL",
        "• IGT: 140-199 mg/dL (2h)",
        "• Lifestyle first-line",
        "• Metformin if BMI ≥35",
        "• Annual monitoring"
    ]

    for point in secondary_points:
        p = secondary_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(10)

    # Tertiary prevention
    tertiary_box = slide.shapes.add_textbox(Inches(6.3), Inches(1.2), Inches(2.8), Inches(2.5))
    tertiary_tf = tertiary_box.text_frame
    tertiary_tf.text = "TERTIARY PREVENTION\n(Prevent Complications)"
    tertiary_tf.paragraphs[0].font.bold = True
    tertiary_tf.paragraphs[0].font.size = Pt(14)
    tertiary_tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)

    tertiary_points = [
        "• Established diabetes",
        "• Optimal glycemic control",
        "• Complication screening",
        "• CV risk management",
        "• Multidisciplinary care",
        "• Lifestyle reinforcement"
    ]

    for point in tertiary_points:
        p = tertiary_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(10)

    # Evidence statement
    evidence_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    evidence_tf = evidence_box.text_frame
    evidence_tf.text = "EVIDENCE: DPP trial showed 58% diabetes prevention with lifestyle alone. Da Qing study: 46% reduction in China. Indian DPP study: 28.5% reduction with lifestyle intervention."
    evidence_tf.paragraphs[0].font.size = Pt(12)
    evidence_tf.paragraphs[0].font.bold = True
    evidence_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_prevention_flowchart_slide(prs):
    """Prevention flowchart visual"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Diabetes Prevention Decision Flowchart"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add prevention flowchart diagram
    try:
        diagram_path = "TLM_Diabetes_Mellitus/visualizations/prevention_flowchart.png"
        slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    except:
        # Fallback flowchart description
        flow_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
        flow_tf = flow_box.text_frame
        flow_tf.text = "PREVENTION DECISION ALGORITHM:\n\n1. RISK ASSESSMENT\n   • IDRS >60 or prediabetes criteria met\n\n2. PRIMARY PREVENTION (High-risk individuals)\n   • Intensive lifestyle intervention (DPP model)\n   • 16 sessions in 24 weeks\n   • 5-7% weight loss + 150 min/week activity\n   • Metformin if BMI ≥35 kg/m²\n\n3. SECONDARY PREVENTION (Prediabetes)\n   • Lifestyle modification\n   • Metformin if indicated\n   • Annual glucose monitoring\n\n4. MONITORING & FOLLOW-UP\n   • Quarterly metabolic assessments\n   • Reinforce lifestyle changes\n   • Address barriers to adherence"
        flow_tf.paragraphs[0].font.size = Pt(12)
        flow_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_national_program_slide(prs):
    """National program overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "National Programme for Prevention & Control of Diabetes"
    title_tf.paragraphs[0].font.size = Pt(28)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.3))
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.text = "NPCDCS - Ministry of Health & Family Welfare, India (2010-ongoing)"
    subtitle_tf.paragraphs[0].font.size = Pt(14)
    subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(127, 140, 141)

    # Key objectives
    objectives_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(3.5), Inches(2))
    objectives_tf = objectives_box.text_frame
    objectives_tf.text = "KEY OBJECTIVES"
    objectives_tf.paragraphs[0].font.bold = True
    objectives_tf.paragraphs[0].font.size = Pt(14)
    objectives_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    objectives = [
        "• 25% reduction in premature mortality",
        "• Early diagnosis and treatment",
        "• Health promotion and prevention",
        "• Capacity building of healthcare workers",
        "• Improved quality of life for patients"
    ]

    for obj in objectives:
        p = objectives_tf.add_paragraph()
        p.text = obj
        p.font.size = Pt(11)

    # Components
    components_box = slide.shapes.add_textbox(Inches(4.5), Inches(1.3), Inches(3.5), Inches(2))
    components_tf = components_box.text_frame
    components_tf.text = "PROGRAM COMPONENTS"
    components_tf.paragraphs[0].font.bold = True
    components_tf.paragraphs[0].font.size = Pt(14)
    components_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    components = [
        "• NCD clinics at district level",
        "• Population-based screening",
        "• Free diagnostic services",
        "• Drug procurement and distribution",
        "• Health promotion activities",
        "• Training and IEC activities"
    ]

    for comp in components:
        p = components_tf.add_paragraph()
        p.text = comp
        p.font.size = Pt(11)

    # Implementation and impact
    impact_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.8))
    impact_tf = impact_box.text_frame
    impact_tf.text = "IMPLEMENTATION: Multi-tier approach (National-State-District-Community). Key achievements: 355 operational NCD clinics, over 4 crore screenings, oral drug distribution to 50 lakh patients. Integration with Ayushman Bharat and health wellness centres."
    impact_tf.paragraphs[0].font.size = Pt(12)
    impact_tf.paragraphs[0].font.bold = True
    impact_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_national_program_visual_slide(prs):
    """National program visual diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "NPCDCS Implementation Framework"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add national program diagram
    try:
        diagram_path = "TLM_Diabetes_Mellitus/visualizations/national_program_diagram.png"
        slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    except:
        # Fallback implementation framework
        framework_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
        framework_tf = framework_box.text_frame
        framework_tf.text = "NPCDCS IMPLEMENTATION FRAMEWORK:\n\nNATIONAL LEVEL\n• Policy formulation, resource allocation\n• Technical guidelines development\n• Monitoring and evaluation\n\nSTATE LEVEL\n• State programme implementation\n• Human resource development\n• Quality assurance\n\nDISTRICT LEVEL\n• NCD clinics establishment\n• Screening camps organization\n• Referral systems setup\n\nCOMMUNITY LEVEL\n• ASHA/ANM worker training\n• Village health and nutrition days\n• Health education campaigns\n\nINTEGRATED APPROACH\n• Screening → Diagnosis → Treatment → Follow-up\n• Free drugs and diagnostics\n• Community participation"
        framework_tf.paragraphs[0].font.size = Pt(12)
        framework_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_future_directions_slide(prs):
    """Future directions and innovations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Future Directions in Diabetes Care (2024-2030)"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Precision medicine
    precision_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3), Inches(2.5))
    precision_tf = precision_box.text_frame
    precision_tf.text = "PRECISION MEDICINE"
    precision_tf.paragraphs[0].font.bold = True
    precision_tf.paragraphs[0].font.size = Pt(15)
    precision_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    precision_points = [
        "• Genetic profiling for therapy",
        "• Biomarker development",
        "• Targeted therapies",
        "• Individualized targets",
        "• Pharmacogenomics"
    ]

    for point in precision_points:
        p = precision_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

    # Advanced technologies
    tech_box = slide.shapes.add_textbox(Inches(3.8), Inches(1.2), Inches(3), Inches(2.5))
    tech_tf = tech_box.text_frame
    tech_tf.text = "ADVANCED TECHNOLOGIES"
    tech_tf.paragraphs[0].font.bold = True
    tech_tf.paragraphs[0].font.size = Pt(15)
    tech_tf.paragraphs[0].font.color.rgb = RGBColor(155, 89, 182)

    tech_points = [
        "• Closed-loop systems",
        "• Stem cell therapy",
        "• Gene therapy (CRISPR)",
        "• AI-driven predictions",
        "• Digital therapeutics"
    ]

    for point in tech_points:
        p = tech_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

    # Health system changes
    system_box = slide.shapes.add_textbox(Inches(7), Inches(1.2), Inches(2.75), Inches(2.5))
    system_tf = system_box.text_frame
    system_tf.text = "HEALTH SYSTEM TRANSFORMATION"
    system_tf.paragraphs[0].font.bold = True
    system_tf.paragraphs[0].font.size = Pt(14)
    system_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    system_points = [
        "• Multidisciplinary teams",
        "• Digital health integration",
        "• Population health approach",
        "• Health economics focus",
        "• Community-based care"
    ]

    for point in system_points:
        p = system_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

    # Challenges and opportunities
    challenges_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    challenges_tf = challenges_box.text_frame
    challenges_tf.text = "CHALLENGES: Cost of new technologies, equitable access, workforce training, data privacy. OPPORTUNITIES: Prevention focus shift, AI integration, global collaborations, personalized approaches."
    challenges_tf.paragraphs[0].font.size = Pt(12)
    challenges_tf.paragraphs[0].font.bold = True
    challenges_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_challenges_solutions_slide(prs):
    """Challenges and strategic solutions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Challenges in Diabetes Control & Prevention"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Key challenges
    challenges_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(3))
    challenges_tf = challenges_box.text_frame
    challenges_tf.text = "KEY CHALLENGES"
    challenges_tf.paragraphs[0].font.bold = True
    challenges_tf.paragraphs[0].font.size = Pt(18)
    challenges_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

    challenges = [
        "• Therapeutic inertia (delayed treatment intensification)",
        "• Long-term lifestyle adherence issues",
        "• Hypoglycemia risk with intensive control",
        "• Resource constraints in LMICs",
        "• Cultural and behavioral barriers",
        "• Socioeconomic disparities in access to care"
    ]

    for challenge in challenges:
        p = challenges_tf.add_paragraph()
        p.text = challenge
        p.font.size = Pt(12)

    # Strategic solutions
    solutions_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(3))
    solutions_tf = solutions_box.text_frame
    solutions_tf.text = "STRATEGIC SOLUTIONS"
    solutions_tf.paragraphs[0].font.bold = True
    solutions_tf.paragraphs[0].font.size = Pt(18)
    solutions_tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)

    solutions = [
        "• Treatment algorithm audits and standardization",
        "• Behavioral support and digital health tools",
        "• CGM-guided therapy adjustments",
        "• Task shifting to trained non-physicians",
        "• Culturally adapted health education",
        "• Community-based peer support programs"
    ]

    for solution in solutions:
        p = solutions_tf.add_paragraph()
        p.text = solution
        p.font.size = Pt(12)

    # Success metrics
    metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1))
    metrics_tf = metrics_box.text_frame
    metrics_tf.text = "MEASURING SUCCESS: Process indicators (screening rates, timely intensification), outcome indicators (HbA1c targets, complication rates), patient-reported outcomes"
    metrics_tf.paragraphs[0].font.size = Pt(11)
    metrics_tf.paragraphs[0].font.bold = True
    metrics_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_conclusion_slide(prs):
    """Comprehensive conclusion and key takeaways"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Conclusion & Key Takeaways"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Key takeaways organized by section
    takeaways_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
    takeaways_tf = takeaways_box.text_frame
    takeaways_tf.text = "EPIDEMIOLOGY & BURDEN"
    takeaways_tf.paragraphs[0].font.bold = True
    takeaways_tf.paragraphs[0].font.size = Pt(16)
    takeaways_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    epi_points = [
        "• 537M adults with diabetes globally, 101M in India",
        "• Projected growth to 783M (2045) and 160M (India)",
        "• South Asia faces highest proportional increase",
        "• LMICs bear disproportionate burden",
        "• Economic cost: $966B globally, INR 2.5-3 lakh crores in India"
    ]

    for point in epi_points:
        p = takeaways_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

    patho_title = takeaways_tf.add_paragraph()
    patho_title.text = "PATHOGENESIS & CLASSIFICATION"
    patho_title.font.bold = True
    patho_title.font.size = Pt(16)
    patho_title.font.color.rgb = RGBColor(46, 204, 113)

    patho_points = [
        "• T2DM: Insulin resistance + β-cell dysfunction",
        "• T1DM: Autoimmune destruction of β-cells",
        "• Prediabetes affects 541M globally",
        "• Genetic-environment interaction crucial",
        "• Early intervention prevents progression"
    ]

    for point in patho_points:
        p = takeaways_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

    control_title = takeaways_tf.add_paragraph()
    control_title.text = "MANAGEMENT & PREVENTION"
    control_title.font.bold = True
    control_title.font.size = Pt(16)
    control_title.font.color.rgb = RGBColor(155, 89, 182)

    control_points = [
        "• Individualized targets: HbA1c <7% generally",
        "• Metformin first-line, early combination therapy",
        "• CGM/TIR important for personalized control",
        "• DPP lifestyle intervention: 58% prevention",
        "• Multidisciplinary, technology-integrated approach"
    ]

    for point in control_points:
        p = takeaways_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

    future_title = takeaways_tf.add_paragraph()
    future_title.text = "FUTURE DIRECTIONS"
    future_title.font.bold = True
    future_title.font.size = Pt(16)
    future_title.font.color.rgb = RGBColor(44, 62, 80)

    future_points = [
        "• Precision medicine and AI-driven care",
        "• Stem cell and gene therapies",
        "• Closed-loop automated systems",
        "• Population health and prevention focus",
        "• Global collaboration for equitable care"
    ]

    for point in future_points:
        p = takeaways_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(11)

    # Call to action
    action_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
    action_tf = action_box.text_frame
    action_tf.text = "CALL TO ACTION: Early screening, intensive lifestyle intervention, comprehensive management, technology adoption, health system strengthening"
    action_tf.paragraphs[0].font.size = Pt(12)
    action_tf.paragraphs[0].font.bold = True
    action_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

def main():
    """Main function to create the comprehensive PowerPoint presentation"""
    # Create presentation
    prs = Presentation()

    # Set slide size to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("Creating comprehensive diabetes presentation...")

    # Create slides
    print("• Title slide...")
    create_title_slide(prs)

    print("• Overview slide...")
    create_overview_slide(prs)

    print("• Definition & classification slides...")
    create_definitions_slide(prs)
    create_types_diabetes_slide(prs)

    print("• Epidemiology slides...")
    create_epidemiology_global_slide(prs)
    create_epidemiology_india_slide(prs)
    create_epidemiology_visual_slide(prs)

    print("• Pathophysiology slides...")
    create_pathophysiology_slide(prs)
    create_pathophysiology_visual_slide(prs)

    print("• Risk factors & clinical features...")
    create_risk_factors_slide(prs)
    create_risk_factors_visual_slide(prs)
    create_clinical_features_slide(prs)

    print("• Diagnosis slide...")
    create_diagnosis_criteria_slide(prs)

    print("• Management slides...")
    create_treatment_overview_slide(prs)
    create_treatment_algorithm_slide(prs)

    print("• Control & monitoring...")
    create_control_targets_slide(prs)
    create_monitoring_slide(prs)

    print("• Prevention slides...")
    create_prevention_overview_slide(prs)
    create_prevention_flowchart_slide(prs)

    print("• National programs...")
    create_national_program_slide(prs)
    create_national_program_visual_slide(prs)

    print("• Future directions & challenges...")
    create_future_directions_slide(prs)
    create_challenges_solutions_slide(prs)

    print("• Conclusion...")
    create_conclusion_slide(prs)

    # Save presentation
    output_path = "TLM_Diabetes_Mellitus/Diabetes_Comprehensive_TLM_Presentation.pptx"
    prs.save(output_path)

    print(f"\n✅ Comprehensive PowerPoint presentation created successfully!")
    print(f"📁 Location: {output_path}")
    print(f"📊 Slides created: {len(prs.slides)}")
    print("\n📋 Presentation covers:")
    print("• Complete epidemiology (global, India, trends)")
    print("• Pathophysiology & molecular mechanisms")
    print("• All diabetes types & classifications")
    print("• Risk factors & clinical presentations")
    print("• Diagnosis criteria & methods")
    print("• Management algorithms & therapies")
    print("• Control strategies & monitoring")
    print("• Prevention frameworks (primary/secondary/tertiary)")
    print("• National programmes (NPCDCS)")
    print("• Future directions & innovations")
    print("• Rich visuals & proper formatting")

if __name__ == "__main__":
    main()
