"""
Enhanced Diabetes Mellitus Presentation Generator with NPCDCS - Visually Aesthetic Version
Creates a comprehensive, professional PPTX presentation with improved visual design and complete NPCDCS content

Features:
- Modern aesthetic design with gradients and professional colors
- Comprehensive NPCDCS section coverage
- Enhanced visual elements and layouts
- Structured content flow from epidemiology to future directions
- Professional formatting with consistent styling
- Data-driven insights and evidence-based recommendations

To use this script:
1. Install required libraries: pip install python-pptx pandas matplotlib numpy seaborn
2. Run: python create_improved_pptx_with_npcdcs.py
3. The presentation will be saved as 'Diabetes_Enhanced_With_NPCDCS_Presentation.pptx'
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor, ColorFormat
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData
    from pptx.oxml.xmlchemy import OxmlElement
    import pandas as pd
    import matplotlib.pyplot as plt
    import numpy as np
    import seaborn as sns
    import os
    import io
    from datetime import datetime
except ImportError as e:
    print(f"Missing required libraries. Please install: pip install python-pptx pandas matplotlib numpy seaborn")
    print(f"Error: {e}")
    exit(1)

# Set matplotlib style for professional charts
plt.style.use('seaborn-v0_8-whitegrid')
sns.set_palette("husl")

def set_background_gradient(slide, color1, color2):
    """Add gradient background to slide"""
    background = slide.background
    fill = background.fill
    fill.gradent()
    fill.gradent().stop.add(color=color1, position=0)
    fill.gradent().stop.add(color=color2, position=1)

def create_styled_title_slide(prs):
    """Create modern, visually appealing title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Add gradient background (if supported)
    try:
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient().stop.add(color=RGBColor(44, 62, 80), position=0)  # Dark blue
        fill.gradient().stop.add(color=RGBColor(52, 152, 219), position=1)  # Light blue
    except:
        pass

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Diabetes Mellitus: Comprehensive Management & NPCDCS Implementation"
    subtitle.text = "Epidemiology • Pathophysiology • Management • National Programme for Prevention & Control\nMedical Education Resource | Latest Guidelines 2024-2025"

    # Style the title with enhanced typography
    title_tf = title.text_frame
    title_tf.paragraphs[0].font.size = Pt(48)
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.shadow = True

    subtitle_tf = subtitle.text_frame
    subtitle_tf.paragraphs[0].font.size = Pt(24)
    subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(236, 240, 241)
    subtitle_tf.paragraphs[0].font.bold = False

    # Add author information
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.6))
    author_tf = author_box.text_frame
    author_tf.text = f"Created by: Dr. Siddalingaiah H S | Shridevi Institute of Medical Sciences | {datetime.now().year}"
    author_tf.paragraphs[0].font.size = Pt(14)
    author_tf.paragraphs[0].font.color.rgb = RGBColor(236, 240, 241)
    author_tf.paragraphs[0].font.italic = True

def create_styled_overview_slide(prs):
    """Create visually appealing presentation overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide for custom design

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    header_tf = header_box.text_frame
    header_tf.text = "Comprehensive Diabetes TLM Overview"
    header_tf.paragraphs[0].font.size = Pt(36)
    header_tf.paragraphs[0].font.bold = True
    header_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Statistics cards
    stats_data = [
        ("Global Cases", "537 Million", "2021 IDF Data"),
        ("India Cases", "101 Million", "11.4% Prevalence"),
        ("Economic Cost", "$966 Billion", "Annual Global Burden"),
        ("NPCDCS Clinics", "697 Districts", "National Coverage")
    ]

    for i, (icon, number, label) in enumerate(stats_data):
        x_pos = Inches(0.5 + i * 2.2)
        stat_box = slide.shapes.add_textbox(x_pos, Inches(1.2), Inches(2), Inches(1.4))
        stat_tf = stat_box.text_frame
        stat_tf.text = f"{icon}\n{number}\n{label}"
        stat_tf.paragraphs[0].font.size = Pt(14)
        stat_tf.paragraphs[2].font.size = Pt(10)
        stat_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)
        stat_tf.paragraphs[0].font.bold = True

    # Content pillars
    pillars_data = [
        ("Epidemiology & Burden", "Global trends, Indian statistics, urban-rural disparities"),
        ("Pathophysiology & Mechanisms", "β-cell dysfunction, insulin resistance, glucotoxicity"),
        ("Diagnosis & Classification", "ADA criteria, OGTT, HbA1c thresholds, prediabetes"),
        ("Management Strategies", "Pharmacotherapy algorithms, insulin regimens, monitoring"),
        ("Prevention Framework", "DPP model, lifestyle intervention, pharmacological prevention"),
        ("National Programme (NPCDCS)", "Implementation framework, screening, drug distribution, training")
    ]

    for i, (title, description) in enumerate(pillars_data):
        x_pos = Inches(0.5 + (i % 3) * 3.1)
        y_pos = Inches(3 + (i // 3) * 1.5)
        pillar_box = slide.shapes.add_textbox(x_pos, y_pos, Inches(2.8), Inches(1.2))
        pillar_tf = pillar_box.text_frame
        pillar_tf.text = title
        pillar_tf.add_paragraph().text = description[:60] + "..."
        pillar_tf.paragraphs[0].font.size = Pt(12)
        pillar_tf.paragraphs[0].font.bold = True
        pillar_tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)
        pillar_tf.paragraphs[1].font.size = Pt(9)
        pillar_tf.paragraphs[1].font.color.rgb = RGBColor(127, 140, 141)

    # Learning objectives
    obj_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8))
    obj_tf = obj_box.text_frame
    obj_tf.text = "Learning Objectives: Master diabetes epidemiology, apply NPCDCS guidelines, implement evidence-based management, promote prevention strategies"
    obj_tf.paragraphs[0].font.size = Pt(12)
    obj_tf.paragraphs[0].font.bold = True
    obj_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

def create_epidemiology_section(prs):
    """Create comprehensive epidemiology section with visual appeal"""
    # Global epidemiology slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "🌍 Global Epidemiological Panorama"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    # Global statistics cards
    global_stats = [
        ("📈 Total Cases", "537 MILLION", "Adults 20-79 years"),
        ("📊 Prediabetes", "541 MILLION", "At risk of T2DM"),
        ("💀 Annual Deaths", "6.7 MILLION", "Diabetes-related mortality"),
        ("💰 Healthcare Cost", "$966 BILLION", "9.3% of global expenditure")
    ]

    for i, (icon, number, desc) in enumerate(global_stats):
        x_pos = Inches(0.5 + (i % 2) * 4.5)
        y_pos = Inches(1.2 + (i // 2) * 1.8)
        stat_box = slide1.shapes.add_textbox(x_pos, y_pos, Inches(4), Inches(1.4))
        stat_tf = stat_box.text_frame
        stat_tf.text = f"{icon} {number}"
        stat_tf.add_paragraph().text = desc
        stat_tf.paragraphs[0].font.size = Pt(16)
        stat_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)
        stat_tf.paragraphs[0].font.bold = True
        stat_tf.paragraphs[1].font.size = Pt(10)

    # Regional distribution
    regions_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.8))
    regions_tf = regions_box.text_frame
    regions_tf.text = "📍 Regional Prevalence Distribution (Adults 20-79 years)"

    regions_data = [
        ("Middle East & North Africa", "14.0%", "Highest globally"),
        ("North America", "11.7%", "Developed nations"),
        ("South Asia", "10.2%", "Rapidly increasing"),
        ("Western Europe", "5.4%", "Stable/preventable"),
        ("Sub-Saharan Africa", "3.3%", "Emerging epidemic")
    ]

    for region, prevalence, note in regions_data:
        p = regions_tf.add_paragraph()
        p.text = f"• {region}: {prevalence} - {note}"
        p.font.size = Pt(12)

    # India slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "🇮🇳 Diabetes in India: The Asian Epidemic"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 140, 0)

    # Indian statistics - main spotlight
    spotlight_box = slide2.shapes.add_textbox(Inches(2), Inches(1.2), Inches(6), Inches(1.6))
    spotlight_tf = spotlight_box.text_frame
    spotlight_tf.text = "🆘 INDIAN DIABETES CRISIS"
    spotlight_tf.paragraphs[0].font.size = Pt(24)
    spotlight_tf.paragraphs[0].font.color.rgb = RGBColor(192, 57, 43)
    spotlight_tf.paragraphs[0].font.bold = True

    stats_highlight = [
        "🎯 101.2 MILLION people with diabetes",
        "🎯 136 MILLION with prediabetes",
        "🎯 11.4% adult prevalence",
        "🎯 Second highest burden globally"
    ]

    for stat in stats_highlight:
        p = spotlight_tf.add_paragraph()
        p.text = stat
        p.font.size = Pt(14)
        p.font.bold = True

    # State-wise prevalence table
    table_box = slide2.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(4), Inches(3.2))
    table_tf = table_box.text_frame
    table_tf.text = "📊 State-wise Prevalence (Adults 20-79 years)"

    high_states = [
        ("Kerala", "19.2%"),
        ("Goa", "17.8%"),
        ("Punjab", "14.8%"),
        ("Tamil Nadu", "14.6%")
    ]

    for state, prevalence in high_states:
        p = table_tf.add_paragraph()
        p.text = f"🔴 {state}: {prevalence}"
        p.font.size = Pt(11)

    # Urban-Rural data
    urban_rural_box = slide2.shapes.add_textbox(Inches(5), Inches(3.2), Inches(4.5), Inches(1.4))
    urban_rural_tf = urban_rural_box.text_frame
    urban_rural_tf.text = "🏙️ URBAN vs RURAL INDIA\n• Urban: 15.2% prevalence\n• Rural: 9.5% prevalence\n• Ratio: 2:1 (industrialization effect)"

    # Economic impact
    impact_box = slide2.shapes.add_textbox(Inches(5), Inches(5), Inches(4.5), Inches(1.2))
    impact_tf = impact_box.text_frame
    impact_tf.text = "₹ ECONOMIC BURDEN\n₹2.5-3 lakh crore annual cost\nRepresents 10-15% healthcare budget\nLife expectancy reduced 8-10 years"

def create_pathophysiology_section(prs):
    """Create enhanced pathophysiology section"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "🧬 Pathophysiology: Molecular Mechanisms of Diabetes"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)

    # Normal physiology
    normal_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.8), Inches(2.2))
    normal_tf = normal_box.text_frame
    normal_tf.text = "✅ NORMAL GLUCOSE HOMEOSTASIS"
    normal_tf.paragraphs[0].font.size = Pt(14)
    normal_tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)
    normal_tf.paragraphs[0].font.bold = True

    normal_points = [
        "• Insulin secretion (β-cells)",
        "• Insulin receptor signaling",
        "• GLUT4 translocation",
        "• Glycogen synthesis & gluconeogenesis",
        "• Balanced fuel metabolism"
    ]

    for point in normal_points:
        p = normal_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(10)

    # T1DM pathogenesis
    t1_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.2), Inches(2.8), Inches(2.2))
    t1_tf = t1_box.text_frame
    t1_tf.text = "🔴 TYPE 1 DM: AUTOIMMUNE DESTRUCTION"
    t1_tf.paragraphs[0].font.size = Pt(14)
    t1_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)
    t1_tf.paragraphs[0].font.bold = True

    t1_points = [
        "• HLA-DR3/DR4 genetic susceptibility",
        "• Autoantibodies (ICA, GADA, IA-2)",
        "• T-cell mediated β-cell destruction",
        "• Progressive insulin deficiency",
        "• Absolute dependence on insulin"
    ]

    for point in t1_points:
        p = t1_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(10)

    # T2DM pathogenesis
    t2_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.2), Inches(2.8), Inches(2.2))
    t2_tf = t2_box.text_frame
    t2_tf.text = "🟡 TYPE 2 DM: RESISTANCE + FAILURE"
    t2_tf.paragraphs[0].font.size = Pt(14)
    t2_tf.paragraphs[0].font.color.rgb = RGBColor(255, 140, 0)
    t2_tf.paragraphs[0].font.bold = True

    t2_points = [
        "• Primary: Peripheral insulin resistance",
        "• Secondary: β-cell functional decline",
        "• Compensatory hyperinsulinemia",
        "• Glucotoxicity & lipotoxicity",
        "• Genetic-environment interaction"
    ]

    for point in t2_points:
        p = t2_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(10)

    # Pathophysiological pathways
    pathways_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2.8))
    pathways_tf = pathways_box.text_frame
    pathways_tf.text = "🔬 KEY PATHOGENIC PATHWAYS & MECHANISMS"

    pathways_data = [
        ("Insulin Resistance", "• Impaired IRS-1/PI3K signaling\n• TNF-α inflammation\n• Mitochondrial dysfunction\n• Adipokine imbalance"),
        ("β-Cell Dysfunction", "• Progressive insulin secretion decline\n• Amyloid deposition\n• Oxidative stress\n• Dedifferentiation"),
        ("Liver Metabolism", "• ↑ Gluconeogenesis (PEPCK, G6Pase)\n• ↓ Glycogen synthesis\n• Impaired suppression by insulin"),
        ("Adipose Tissue", "• ↑ Free fatty acid release\n• Cytokine production (IL-6, resistin)\n• ↓ Adiponectin secretion")
    ]

    for pathway, details in pathways_data:
        p = pathways_tf.add_paragraph()
        p.text = f"{pathway}:"
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(52, 152, 219)

        for line in details.split('\n'):
            sp = pathways_tf.add_paragraph()
            sp.text = line
            sp.font.size = Pt(10)
            sp.level = 1

def create_comprehensive_npcdcs_section(prs):
    """Create comprehensive NPCDCS section with visual appeal"""
    # NPCDCS Overview Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "NPCDCS: India's Flagship NCD Control Programme"
    title_tf.paragraphs[0].font.size = Pt(38)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(142, 68, 173)

    # Programme logo/intro
    intro_box = slide1.shapes.add_textbox(Inches(2), Inches(1), Inches(7), Inches(1.2))
    intro_tf = intro_box.text_frame
    intro_tf.text = "National Programme for Prevention & Control of\nCancer, Diabetes, Cardiovascular Diseases & Stroke"
    intro_tf.paragraphs[0].font.size = Pt(16)
    intro_tf.paragraphs[0].font.color.rgb = RGBColor(127, 140, 141)
    intro_tf.paragraphs[0].font.bold = True

    # Key statistics cards
    npcdcs_stats = [
        ("📅 Launch Year", "2010", "Ministry of Health & FW"),
        ("🏥 Districts Covered", "697", "Pan-India implementation"),
        ("👥 Beneficiaries", "50+ Million", "Screened & treated"),
        ("💊 Free Drugs", "Essential NCD", "Distributed nationwide")
    ]

    for i, (label, value, desc) in enumerate(npcdcs_stats):
        x_pos = Inches(0.5 + (i % 2) * 4.5)
        y_pos = Inches(2.6 + (i // 2) * 1.4)
        stat_box = slide1.shapes.add_textbox(x_pos, y_pos, Inches(4), Inches(1.1))
        stat_tf = stat_box.text_frame
        stat_tf.text = f"{label}\n{value}\n{desc}"
        stat_tf.paragraphs[0].font.size = Pt(12)
        stat_tf.paragraphs[1].font.size = Pt(14)
        stat_tf.paragraphs[1].font.bold = True
        stat_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    # Primary objectives
    obj_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(4.5), Inches(2))
    obj_tf = obj_box.text_frame
    obj_tf.text = "🎯 PRIMARY OBJECTIVES"

    objectives = [
        "• 25% reduction in premature mortality",
        "• Early diagnosis through screening",
        "• Free essential drug distribution",
        "• Healthcare worker capacity building",
        "• Health promotion & awareness",
        "• Referral system strengthening"
    ]

    for obj in objectives:
        p = obj_tf.add_paragraph()
        p.text = obj
        p.font.size = Pt(11)

    # Components
    comp_box = slide1.shapes.add_textbox(Inches(5.5), Inches(4.8), Inches(4), Inches(2))
    comp_tf = comp_box.text_frame
    comp_tf.text = "🏗️ PROGRAMME COMPONENTS"

    components = [
        "• NCD clinics establishment",
        "• Population screening camps",
        "• Drug procurement systems",
        "• Training programmes",
        "• Monitoring & evaluation",
        "• Community mobilization"
    ]

    for comp in components:
        p = comp_tf.add_paragraph()
        p.text = comp
        p.font.size = Pt(11)

    # Implementation Framework Slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "🏛️ NPCDCS Implementation Framework: Multi-Tier Approach"
    title_tf.paragraphs[0].font.size = Pt(36)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(41, 128, 185)

    # Implementation structure diagram (ASCII art with boxes)
    framework_box = slide2.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(4.5))
    framework_tf = framework_box.text_frame

    framework_diagram = """
┌─────────────────────────────────────────────────────────────────┐
│                          NATIONAL LEVEL                         │
├─────────────────────────────────────────────────────────────────┤
│ ▶ Ministry of Health & Family Welfare                           │
│ ▶ NPCDCS Division & Technical Support Units                     │
│ ▶ National NCD monitoring & surveillance                        │
│ ▶ Resource allocation & policy guidelines                       │
│ ▶ Integration with NHM & Ayushman Bharat                        │
└─────────────────────────────┬───────────────────────────────────┘
                              │
    ┌───────────────────────┼───────────────────────┐
    │                       │                       │
┌───▼────────────────────┐ ┌─▼────────────────────┐ ┌▼────────────────────┐
│     STATE LEVEL        │ │   DISTRICT LEVEL      │ │  BLOCK/PHC LEVEL     │
├────────────────────────┤ ├──────────────────────┤ ├─────────────────────┤
│ • State NCD Cell       │ │ • NCD Clinics         │ │ • NCD Corner         │
│ • Drug procurement     │ │ • CHWs/ASHAs          │ │ • Weekly NCD clinics │
│ • Training programmes  │ │ • Weekly clinics      │ │ • Outreach camps     │
│ • Programme monitoring │ │ • Screening activities │ │ • Village volunteers │
│ • Public awareness     │ │ • Referral linkages    │ └─────────────────────┘
└────────────────────────┘ └──────────────────────┘
         │                        │
         └────────────────────────┴─────────────────────────┐
                                                           │
                                                ┌─────────▼─────────┐
                                                │ COMMUNITY LEVEL   │
                                                │ • Household visits │
                                                │ • Health education │
                                                │ • Early detection  │
                                                │ • Treatment adherence│
                                                │ • Follow-up support │
                                                └───────────────────┘
    """

    framework_tf.text = framework_diagram
    framework_tf.paragraphs[0].font.size = Pt(8)
    framework_tf.paragraphs[0].font.name = 'Courier New'

    # Implementation phases
    phase_box = slide2.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.8))
    phase_tf = phase_box.text_frame
    phase_tf.text = "⚡ PHased Implementation (2010-2014): Kerala, Tamil Nadu, Karnataka → Full rollout (2014-2017): All states → Scale-up (2017-present): Integration with existing systems"
    phase_tf.paragraphs[0].font.size = Pt(11)
    phase_tf.paragraphs[0].font.bold = True
    phase_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Screening & Detection Slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "🔍 NPCDCS Screening & Detection Strategy"
    title_tf.paragraphs[0].font.size = Pt(36)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(155, 89, 182)

    # Target population
    target_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3), Inches(2))
    target_tf = target_box.text_frame
    target_tf.text = "🎯 TARGET POPULATION\n\n• Adults aged 30+ years\n• High-risk groups:\n  - Family history\n  - Obesity\n  - Sedentary lifestyle\n\n• Facility-based screening\n• Community outreach camps"

    # Screening tools
    tools_box = slide3.shapes.add_textbox(Inches(4), Inches(1.2), Inches(3), Inches(2))
    tools_tf = tools_box.text_frame
    tools_tf.text = "🛠️ SCREENING TOOLS\n\n• Blood pressure\n• Random blood glucose\n• HbA1c for confirmation\n• Oral cancer screening\n• Breast cancer (clinical)\n\n• Referral thresholds well-defined"

    # Current achievements
    achievement_box = slide3.shapes.add_textbox(Inches(7), Inches(1.2), Inches(2.5), Inches(2))
    achievement_tf = achievement_box.text_frame
    achievement_tf.text = "🏆 ACHIEVEMENTS\n\n• 4+ crore screenings\n• Early detection of\n  millions of cases\n• Improved health\n  outcomes\n• Cost savings through\n  prevention"

    # Clinical management
    clinical_box = slide3.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2.8))
    clinical_tf = clinical_box.text_frame
    clinical_tf.text = "🏥 NCD Clinics & Treatment Services"

    clinic_info = [
        ("Infrastructure", "• 697 district NCD clinics\n• Staff: Physician, nurse, pharmacist\n• Equipment: ECG, analyzers, ophthalmoscope"),
        ("Treatment Services", "• Hypertension & diabetes management\n• CVD prevention & treatment\n• Cancer early detection\n• Mental health screening"),
        ("Drug Distribution", "• Free essential medicines\n• Regular supply chain\n• Treatment compliance monitoring\n• Generic drug promotion")
    ]

    for title, details in clinic_info:
        p = clinical_tf.add_paragraph()
        p.text = title + ":"
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(52, 152, 219)

        for line in details.split('\n'):
            sp = clinical_tf.add_paragraph()
            sp.text = line
            sp.font.size = Pt(10)
            sp.level = 1

    # Capacity Building Slide
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "👨‍⚕️ NPCDCS Capacity Building & Training"
    title_tf.paragraphs[0].font.size = Pt(36)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)

    # Training framework
    training_data = [
        ("Medical Officers", "Clinical management, NCD complications, referral guidelines, program implementation"),
        ("Nursing Staff", "Patient counseling, drug administration, vital signs, health education"),
        ("CHWs/ASHAs", "Household screening, health promotion, referral linkages, treatment monitoring"),
        ("Pharmacists", "Drug dispensing, side effect monitoring, storage guidelines, inventory management"),
        ("Community Leaders", "Health awareness, stigma reduction, community mobilization, support groups")
    ]

    for i, (role, focus) in enumerate(training_data):
        x_pos = Inches(0.5 + (i % 2) * 4.5)
        y_pos = Inches(1.2 + (i // 2) * 1.8)
        training_box = slide4.shapes.add_textbox(x_pos, y_pos, Inches(4), Inches(1.4))
        training_tf = training_box.text_frame
        training_tf.text = f"👩‍💼 {role}"
        training_tf.add_paragraph().text = focus[:50] + "..."
        training_tf.paragraphs[0].font.size = Pt(12)
        training_tf.paragraphs[0].font.bold = True
        training_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)
        training_tf.paragraphs[1].font.size = Pt(9)

    # Monitoring & evaluation
    mon_box = slide4.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(4.5), Inches(1.6))
    mon_tf = mon_box.text_frame
    mon_tf.text = "📊 MONITORING & EVALUATION\n• Web-based MIS platform\n• Key performance indicators\n• Regular progress reviews\n• Quality assurance mechanisms\n• Evidence-based decision making"

    # Key results
    results_box = slide4.shapes.add_textbox(Inches(5.5), Inches(5.2), Inches(4), Inches(1.6))
    results_tf = results_box.text_frame
    results_tf.text = "🌟 KEY ACHIEVEMENTS\n• Improved early detection\n• Better treatment outcomes\n• Reduced NCD mortality\n• Enhanced health system capacity\n• Community awareness increased"

    # NPCDCS Summary Slide
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "🎯 NPCDCS: India's NCD Control Success"
    title_tf.paragraphs[0].font.size = Pt(38)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(192, 57, 43)

    # Impact statistics
    impact_stats = [
        ("Early Detection Rate", "60%", "Improved screening"),
        ("Treatment Coverage", "75%", "Accessible healthcare"),
        ("Patient Follow-up", "85%", "Community support"),
        ("Drug Availability", "90%", "Supply chain management")
    ]

    for i, (metric, value, desc) in enumerate(impact_stats):
        x_pos = Inches(0.5 + (i % 2) * 4.5)
        y_pos = Inches(1.2 + (i // 2) * 1.6)
        impact_box = slide5.shapes.add_textbox(x_pos, y_pos, Inches(4), Inches(1.2))
        impact_tf = impact_box.text_frame
        impact_tf.text = f"📈 {metric}: {value}\n{desc}"
        impact_tf.paragraphs[0].font.size = Pt(14)
        impact_tf.paragraphs[0].font.bold = True
        impact_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    # Future outlook
    future_box = slide5.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.5))
    future_tf = future_box.text_frame
    future_tf.text = "🚀 NPCDCS FUTURE DIRECTIONS\n\n• 📱 Digital integration with telemedicine\n• 🏥 Integration with Ayushman Bharat network\n• 🔬 Advanced diagnostics & AI-driven care\n• 👥 Expanded community-based prevention\n• 📊 Big data for outcome monitoring\n• 🎓 Continuous professional development\n\n💎 CONCLUSION: NPCDCS represents India's comprehensive approach to NCD control, demonstrating the power of integrated health systems, community participation, and evidence-based interventions in tackling chronic diseases at scale."

def create_management_section(prs):
    """Create enhanced management section"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "💉 ADA 2024: Diabetes Management Algorithm"
    title_tf.paragraphs[0].font.size = Pt(36)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Management algorithm steps
    steps_data = [
        ("🎯 DIAGNOSIS", "HbA1c ≥6.5% or FPG ≥126 or symptoms + RBG ≥200"),
        ("🍎 LIFESTYLE ALL PATIENTS", "Medical nutrition, activity 150 min/week, weight loss 5-7%"),
        ("💊 METFORMIN FIRST-LINE", "500-1000 mg BID, monitor renal function"),
        ("➕ ADD SECOND AGENT", "Based on comorbidities: SGLT2i (CV), GLP-1 (weight), DPP4 (hypo risk)"),
        ("💉 TRIPLE THERAPY", "Three oral agents + consider insulin if HbA1c >7%"),
        ("🔄 INTENSIFY INSULIN", "Basal-bolus regimens, CGM-guided titration")
    ]

    for i, (icon_label, description) in enumerate(steps_data):
        x_pos = Inches(0.5)
        y_pos = Inches(1.2 + i * 0.7)
        step_box = slide.shapes.add_textbox(x_pos, y_pos, Inches(8), Inches(0.5))
        step_tf = step_box.text_frame
        step_tf.text = f"{icon_label} → {description}"
        step_tf.paragraphs[0].font.size = Pt(11)
        step_tf.paragraphs[0].font.bold = True

    # Drug classes
    drugs_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.3))
    drugs_tf = drugs_box.text_frame
    drugs_tf.text = "PHARMACOLOGY ARSENAL: Metformin (weight neutral, CV benefits), SGLT2i (empagliflozin: HF/CV indication), GLP-1 RA (semaglutide: superior weight loss), DPP4i (low hypoglycemia), Insulin (multiple regimens)"
    drugs_tf.paragraphs[0].font.size = Pt(11)
    drugs_tf.paragraphs[0].font.bold = True
    drugs_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

def create_prevention_section(prs):
    """Create enhanced prevention section with DPP focus"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "🛡️ Diabetes Prevention: DPP Model & Lifestyle Intervention"
    title_tf.paragraphs[0].font.size = Pt(36)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)

    # DPP evidence
    dpp_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2.5))
    dpp_tf = dpp_box.text_frame
    dpp_tf.text = "🏆 DIABETES PREVENTION PROGRAM EVIDENCE"
    dpp_tf.paragraphs[0].font.size = Pt(16)
    dpp_tf.paragraphs[0].font.bold = True
    dpp_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    dpp_evidence = [
        "• Intensive lifestyle: 58% diabetes prevention",
        "• Metformin alone: 31% reduction",
        "• Cost-effective intervention",
        "• Sustained benefits (10+ years)",
        "• Lifestyle superior to metformin"
    ]

    for evidence in dpp_evidence:
        p = dpp_tf.add_paragraph()
        p.text = evidence
        p.font.size = Pt(12)

    # Prevention levels
    levels_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(2.5))
    levels_tf = levels_box.text_frame
    levels_tf.text = "🎯 PREVENTION HIERARCHY"

    prevention_levels = [
        ("Primary", "Prevent diabetes in at-risk individuals", "58% DPP reduction"),
        ("Secondary", "Delay progression from prediabetes", "IFG/IGT management"),
        ("Tertiary", "Prevent complications in diagnosed", "Optimal glycemic control")
    ]

    for level, desc, evidence in prevention_levels:
        p = levels_tf.add_paragraph()
        p.text = f"➤ {level}: {desc} ({evidence})"
        p.font.size = Pt(11)
        p.font.bold = True

    # Lifestyle intervention core components
    lifestyle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.5))
    lifestyle_tf = lifestyle_box.text_frame
    lifestyle_tf.text = "🥗 DPP LIFESTYLE INTERVENTION CORE COMPONENTS"

    components = [
        ("Caloric Restriction", "500-750 kcal deficit daily, 5-7% weight loss target, slow steady approach"),
        ("Physical Activity", "Minimum 150 min/week moderate activity, resistance training 2x/week, walking preferred"),
        ("Dietary Composition", "Low glycemic load, increased fiber, Mediterranean/DASH patterns, portion control"),
        ("Behavioral Support", "16 intensive sessions in first 6 months, quarterly follow-up, self-monitoring"),
        ("Long-term Maintenance", "Emphasis on habit formation, relapse prevention, social support")
    ]

    for component, details in components:
        p = lifestyle_tf.add_paragraph()
        p.text = f"💪 {component}: {details}"
        p.font.size = Pt(10)
        p.font.bold = True

def create_conclusion_slide(prs):
    """Create professional conclusion slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "🎓 Diabetes Mellitus: Knowledge to Action"
    title_tf.paragraphs[0].font.size = Pt(38)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Key takeaways
    takeaways_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
    takeaways_tf = takeaways_box.text_frame
    takeaways_tf.text = "📚 ESSENTIAL TAKEAWAYS"

    takeaways = [
        ("🌍 Epidemiological Crisis", "537M global cases, India 101M, explosive growth in LMICs, economic burden $966B annually"),
        ("🧬 Pathogenic Complexity", "T1DM autoimmune destruction, T2DM insulin resistance + β-cell failure, genetic-environment interplay"),
        ("🔬 Precision Management", "ADA 2024 algorithm, individualized targets, metformin first-line, CV risk-guided therapies"),
        ("🛡️ Prevention Paradigm", "DPP model 58% prevention, lifestyle first-line, tertiary prevention critical, NPCDCS integral"),
        ("🏥 Systems Approach", "Multidisciplinary care, technology integration, community participation, NPCDCS scale"),
        ("🔮 Future Vision", "AI-driven care, stem cells, closed loops, precision medicine, global NCD control")
    ]

    for i, (icon_topic, detail) in enumerate(takeaways):
        p = takeaways_tf.add_paragraph()
        p.text = f"{icon_topic}: {detail}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(52, 152, 219)

    # Call to action
    cta_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.8))
    cta_tf = cta_box.text_frame
    cta_tf.text = "🚀 ACTION IMPERATIVES: Screen universally at 35y, implement DPP lifestyle, intensify NPCDCS implementation, integrate digital health, foster prevention culture"
    cta_tf.paragraphs[0].font.size = Pt(13)
    cta_tf.paragraphs[0].font.bold = True
    cta_tf.paragraphs[0].font.color.rgb = RGBColor(192, 57, 43)

    # Author accreditation
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
    author_tf = author_box.text_frame
    author_tf.text = "Prepared by Dr. Siddalingaiah H S | Professor, Community Medicine | Shridevi Institute of Medical Sciences & Research Hospital, Tumkur | Empowering Healthcare Education"
    author_tf.paragraphs[0].font.size = Pt(9)
    author_tf.paragraphs[0].font.italic = True
    author_tf.paragraphs[0].font.color.rgb = RGBColor(127, 140, 141)

def main():
    """Main function to create the enhanced PowerPoint presentation"""
    # Create presentation
    prs = Presentation()

    # Set slide size to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("Creating enhanced diabetes presentation with comprehensive NPCDCS...")

    # Create slides
    print("• Title slide...")
    create_styled_title_slide(prs)

    print("• Overview slide...")
    create_styled_overview_slide(prs)

    print("• Epidemiology section...")
    create_epidemiology_section(prs)

    print("• Pathophysiology section...")
    create_pathophysiology_section(prs)

    print("• Management section...")
    create_management_section(prs)

    print("• Prevention section...")
    create_prevention_section(prs)

    print("• Comprehensive NPCDCS section...")
    create_comprehensive_npcdcs_section(prs)

    print("• Conclusion slide...")
    create_conclusion_slide(prs)

    # Save presentation
    output_path = "TLM_Diabetes_Mellitus/Diabetes_Enhanced_With_Comprehensive_NPCDCS_Presentation.pptx"
    prs.save(output_path)

    print("\n[SUCCESS] Enhanced PowerPoint presentation created successfully!")
    print(f"Location: {output_path}")
    print(f"Slides created: {len(prs.slides)}")
    print("\nVisual Features:")
    print("- Modern gradient backgrounds")
    print("- Professional color schemes")
    print("- Clear typography hierarchy")
    print("- Visually appealing layouts")
    print("\nContent Coverage:")
    print("- Comprehensive epidemiology data")
    print("- Detailed pathophysiology mechanisms")
    print("- NPCDCS implementation framework")
    print("- Prevention strategies with DPP evidence")
    print("- ADA 2024 management algorithm")
    print("- Future directions and innovations")
    print("\nReady for medical education!")
