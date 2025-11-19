#!/usr/bin/env python3
"""
Simple test for PowerPoint generation
"""

print("Testing PowerPoint Generation...")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    print("Libraries imported successfully")

    # Create simple presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Enhanced Diabetes PPTX Test"
    subtitle.text = "Success!"

    # Save
    prs.save("simple_test.pptx")
    print("Test PowerPoint created successfully")

except Exception as e:
    print(f"Error: {e}")
