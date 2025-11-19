#!/usr/bin/env python3
"""
Clean PowerPoint Creation Script - No Unicode Characters
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    print("Creating clean PowerPoint presentation...")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    print("Adding title slide...")
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Diabetes Mellitus: Comprehensive Management & NPCDCS Implementation"
    subtitle.text = "Epidemiology • Pathophysiology • Management • National Programme"

    # Overview slide
    print("Adding overview slide...")
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(6))
    box.text_frame.text = """DIABETES TLM OVERVIEW

- Epidemiology & Burden: Global trends, Indian statistics, 537M cases worldwide
- Pathophysiology: Type 1 vs Type 2 DM mechanisms
- Management: ADA 2024 algorithm, evidence-based therapies
- Prevention: DPP model, lifestyle intervention (58% effectiveness)
- NPCDCS: India's national programme implementation framework

Created for medical education and training."""

    # Epidemiology slide
    print("Adding epidemiology slide...")
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(0.8))
    box.text_frame.text = "GLOBAL EPIDEMIOLOGY & INDIAN CONTEXT"
    box = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    box.text_frame.text = """
GLOBAL STATISTICS:
- 537 Million people with diabetes (2021)
- 6.7 Million deaths annually
- $966 Billion healthcare cost

INDIAN CRISIS:
- 101.2 Million cases (11.4% prevalence)
- Second highest burden globally
- 136 Million with prediabetes
- Urban-rural ratio: 2:1

HIGH-RISK STATES: Kerala (19.2%), Goa (17.8%), Punjab (14.8%)"""

    # Prevention section
    print("Adding prevention slide...")
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(0.8))
    box.text_frame.text = "DIABETES PREVENTION PROGRAM (DPP) MODEL"
    box = slide4.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    box.text_frame.text = """
EVIDENCE-BASED PREVENTION:
- Intensive lifestyle intervention: 58% diabetes prevention
- Metformin monotherapy: 31% reduction
- Sustained benefits up to 10+ years

NPCDCS INTEGRATION:
- Community-based screening programs
- Free essential medicines distribution
- Health worker training and capacity building
- Nationwide implementation across 697 districts

KEY COMPONENTS:
- Population screening and early detection
- Pharmacological and lifestyle interventions
- Referral systems and treatment monitoring"""

    # Management slide
    print("Adding management slide...")
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(0.8))
    box.text_frame.text = "ADA 2024: DIABETES MANAGEMENT ALGORITHM"
    box = slide5.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    box.text_frame.text = """STEPWISE MANAGEMENT APPROACH:

1. DIAGNOSIS: HbA1c >=6.5% OR FPG >=126 OR symptoms+RBG>=200
2. LIFESTYLE: Medical nutrition, 150 min/week physical activity
3. METFORMIN FIRST-LINE: 500-1000 mg BID, monitor renal function
4. ADD SECOND AGENT: Based on comorbidities
   - CV disease: SGLT2 inhibitors (empagliflozin)
   - Weight management: GLP-1 receptor agonists
   - Hypoglycemia risk: DPP-4 inhibitors
5. TRIPLE THERAPY: Three oral agents + insulin if HbA1c >7%
6. INSULIN INTENSIFICATION: Basal-bolus regimens, CGM-guided

PHARMACOLOGY:
- Metformin (weight neutral, CV benefits)
- SGLT2i (cardiovascular/heart failure indication)
- GLP-1 RA (superior weight loss, CV benefits)
- DPP-4i (low hypoglycemia risk)"""

    # Conclusion slide
    print("Adding conclusion slide...")
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(0.8))
    box.text_frame.text = "DIABETES MELLITUS: KNOWLEDGE TO ACTION"
    box = slide6.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    box.text_frame.text = """
ESSENTIAL TAKEAWAYS:

1. GLOBAL BURDEN: 537M cases, India 101M, explosive growth in LMICs
2. PATHOGENESIS: T1DM autoimmune destruction, T2DM insulin resistance + beta-cell failure
3. MANAGEMENT: ADA 2024 algorithm, individualized targets, metformin first-line
4. PREVENTION: DPP lifestyle model 58% effective, NPCDCS crucial
5. SYSTEMS APPROACH: Multidisciplinary care, technology integration, community participation

ACTION IMPERATIVES:
- Universal screening at age 35+
- DPP lifestyle intervention implementation
- NPCDCS service expansion
- Digital health integration
- Prevention culture development

Created by: Dr. Siddalingaiah H S
Professor, Community Medicine
Shridevi Institute of Medical Sciences"""

    # Save presentation
    output_file = "Diabetes_Enhanced_With_Comprehensive_NPCDCS_Presentation.pptx"
    print(f"Saving presentation as: {output_file}")
    prs.save(output_file)

    # Verify file was created
    import os
    if os.path.exists(output_file):
        file_size = os.path.getsize(output_file)
        print("SUCCESS: Presentation created successfully!")
        print(f"File: {output_file}")
        print(f"Size: {file_size} bytes")
        print(f"Slides: {len(prs.slides)}")
    else:
        print("ERROR: File was not created!")

except Exception as e:
    import traceback
    print(f"ERROR occurred: {e}")
    traceback.print_exc()
