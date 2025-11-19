"""
Diabetes Control, Prevention & National Programme - PowerPoint Presentation Generator
Creates an elaborate PPTX file focused on diabetes control strategies, prevention, and national programmes

To use this script:
1. Install required libraries: pip install python-pptx pandas matplotlib
2. Run: python create_control_prevention_pptx.py
3. The presentation will be saved as 'Diabetes_Control_Prevention_Presentation.pptx'
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    import pandas as pd
    import matplotlib.pyplot as plt
    import numpy as np
    import io
    import os
except ImportError as e:
    print(f"Missing required libraries. Please install: pip install python-pptx pandas matplotlib")
    print(f"Error: {e}")
    exit(1)

def create_title_slide(prs):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Diabetes Mellitus Control, Prevention & National Programme"
    subtitle.text = "Comprehensive Strategies for Diabetes Management in India\nMBBS Teaching Learning Material"

    # Style the title
    title_tf = title.text_frame
    title_tf.paragraphs[0].font.size = Pt(38)
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    title_tf.paragraphs[0].font.bold = True

def create_overview_slide(prs):
    """Create presentation overview slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = "Presentation Overview"
    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = ""

    overview_points = [
        "• Diabetes Control Strategies - Glycemic targets, monitoring, technology",
        "• Prevention Framework - Primary, secondary, and tertiary prevention",
        "• National Programme (NPCDCS) - India's diabetes control initiative",
        "• Implementation Strategies - Multidisciplinary approach",
        "• Evidence-based Interventions - Lifestyle, pharmacological, behavioral",
        "• Outcome Measures - Process and outcome indicators"
    ]

    for point in overview_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.level = 0

def create_control_targets_slide(prs):
    """Glycemic control targets slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Glycemic Control Targets (ADA 2024)"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Left column - General targets
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(3))
    left_tf = left_box.text_frame
    left_tf.text = "GENERAL TARGETS"
    left_tf.paragraphs[0].font.bold = True
    left_tf.paragraphs[0].font.size = Pt(18)
    left_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    general_targets = [
        "• HbA1c: <7.0% (most adults)",
        "• Preprandial glucose: 80-130 mg/dL",
        "• Peak postprandial: <180 mg/dL",
        "• Bedtime glucose: 100-140 mg/dL",
        "• TIR (Time in Range): >70%"
    ]

    for target in general_targets:
        p = left_tf.add_paragraph()
        p.text = target
        p.font.size = Pt(16)

    # Right column - Individualized targets
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(3))
    right_tf = right_box.text_frame
    right_tf.text = "INDIVIDUALIZED TARGETS"
    right_tf.paragraphs[0].font.bold = True
    right_tf.paragraphs[0].font.size = Pt(18)
    right_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    individualized_targets = [
        "• Young patients (<40): HbA1c <6.5-7.0%",
        "• Elderly (≥65): HbA1c 7.5-8.0%",
        "• Longstanding DM (>10y): More liberal",
        "• High CV risk: Stringent control",
        "• Frequent hypoglycemia: Relaxed targets"
    ]

    for target in individualized_targets:
        p = right_tf.add_paragraph()
        p.text = target
        p.font.size = Pt(16)

    # Bottom notes
    notes_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    notes_tf = notes_box.text_frame
    notes_tf.text = "Key Considerations: Age, comorbidities, life expectancy, hypoglycmia risk, and patient preferences should guide target setting. SMBG/CGM data crucial for pattern management."
    notes_tf.paragraphs[0].font.size = Pt(14)
    notes_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_monitoring_slide(prs):
    """Monitoring strategies slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Diabetes Monitoring Strategies"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # SMBG recommendations
    smbg_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(2.5))
    smbg_tf = smbg_box.text_frame
    smbg_tf.text = "SELF-MONITORING OF BLOOD GLUCOSE"
    smbg_tf.paragraphs[0].font.bold = True
    smbg_tf.paragraphs[0].font.size = Pt(16)
    smbg_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    smbg_points = [
        "• Type 2 DM on oral agents: 1-2x daily",
        "• Type 2 DM on insulin: 2-4x daily",
        "• Type 1 DM: 4-6x daily (basal-bolus)",
        "• Pre/postprandial monitoring essential",
        "• Pattern recognition for dose adjustment"
    ]

    for point in smbg_points:
        p = smbg_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)

    # CGM recommendations
    cgm_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(2.5))
    cgm_tf = cgm_box.text_frame
    cgm_tf.text = "CONTINUOUS GLUCOSE MONITORING"
    cgm_tf.paragraphs[0].font.bold = True
    cgm_tf.paragraphs[0].font.size = Pt(16)
    cgm_tf.paragraphs[0].font.color.rgb = RGBColor(155, 89, 182)

    cgm_points = [
        "• TIR (>70%): Primary outcome measure",
        "• TBR (<4%): Minimize hypoglycemia",
        "• TAR (<25%): Manage hyperglycemia",
        "• Ambulatory Glucose Profile analysis",
        "• Real-time vs. intermittently scanned"
    ]

    for point in cgm_points:
        p = cgm_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)

    # HbA1c monitoring
    hba1c_box = slide.shapes.add_textbox(Inches(2.5), Inches(4), Inches(5), Inches(1))
    hba1c_tf = hba1c_box.text_frame
    hba1c_tf.text = "HbA1c: Every 3-6 months (clinically meaningful change: 0.5%)"
    hba1c_tf.paragraphs[0].font.size = Pt(16)
    hba1c_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    hba1c_tf.paragraphs[0].font.bold = True

def create_prevention_levels_slide(prs):
    """Prevention levels slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Levels of Diabetes Prevention"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Primary Prevention
    primary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.8), Inches(4))
    primary_tf = primary_box.text_frame
    primary_tf.text = "PRIMARY PREVENTION\n(Prevent Diabetes Onset)"
    primary_tf.paragraphs[0].font.bold = True
    primary_tf.paragraphs[0].font.size = Pt(18)
    primary_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

    primary_points = [
        "• Target: High-risk individuals (prediabetes)",
        "• DPP model: 58% incidence reduction",
        "• Intensive lifestyle intervention",
        "• 5-7% weight loss + 150 min/week activity",
        "• Metformin in selected cases (BMI ≥35)",
        "• Health promotion campaigns"
    ]

    for point in primary_points:
        p = primary_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(12)

    # Secondary Prevention
    secondary_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.2), Inches(2.8), Inches(4))
    secondary_tf = secondary_box.text_frame
    secondary_tf.text = "SECONDARY PREVENTION\n(Delay Progression)"
    secondary_tf.paragraphs[0].font.bold = True
    secondary_tf.paragraphs[0].font.size = Pt(18)
    secondary_tf.paragraphs[0].font.color.rgb = RGBColor(243, 156, 18)

    secondary_points = [
        "• Target: Prediabetes management",
        "• IFG: 100-125 mg/dL",
        "• IGT: 140-199 mg/dL (2h PG)",
        "• Lifestyle modification first-line",
        "• Metformin optional (BMI ≥35)",
        "• Annual monitoring for progression"
    ]

    for point in secondary_points:
        p = secondary_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(12)

    # Tertiary Prevention
    tertiary_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.2), Inches(2.8), Inches(4))
    tertiary_tf = tertiary_box.text_frame
    tertiary_tf.text = "TERTIARY PREVENTION\n(Prevent Complications)"
    tertiary_tf.paragraphs[0].font.bold = True
    tertiary_tf.paragraphs[0].font.size = Pt(18)
    tertiary_tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)

    tertiary_points = [
        "• Target: Established diabetes",
        "• Optimal glycemic control",
        "• Complication screening (annual)",
        "• Cardiovascular risk management",
        "• Multidisciplinary care approach",
        "• Rehabilitation and support"
    ]

    for point in tertiary_points:
        p = tertiary_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(12)

def create_control_strategies_visual_slide(prs):
    """Control strategies pyramid visual slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Diabetes Control Strategies Pyramid"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add control strategies diagram
    try:
        diagram_path = "TLM_Diabetes_Mellitus/visualizations/control_strategies_diagram.png"
        slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    except:
        # Fallback text
        fallback_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
        fallback_tf = fallback_box.text_frame
        fallback_tf.text = "Control strategies build hierarchically:\n\n1. Glycemic targets & monitoring\n2. Lifestyle interventions\n3. Pharmacological management\n4. Technology integration\n5. Behavioral support\n\nFoundation: Comprehensive multidisciplinary care"
        for para in fallback_tf.paragraphs:
            para.font.size = Pt(16)

def create_national_program_slide(prs):
    """India National Programme for Diabetes Control slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "National Programme for Prevention & Control of Diabetes"
    title_tf.paragraphs[0].font.size = Pt(28)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.3))
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.text = "NPCDCS - Ministry of Health & Family Welfare, India (2010-ongoing)"
    subtitle_tf.paragraphs[0].font.size = Pt(14)
    subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(127, 140, 141)

    # Key objectives
    objectives_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.25), Inches(2.5))
    objectives_tf = objectives_box.text_frame
    objectives_tf.text = "KEY OBJECTIVES"
    objectives_tf.paragraphs[0].font.bold = True
    objectives_tf.paragraphs[0].font.size = Pt(16)
    objectives_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    objectives = [
        "• 25% reduction in premature mortality",
        "• Early diagnosis and treatment",
        "• Health promotion and prevention",
        "• Capacity building of healthcare workers",
        "• Improved quality of life for patients"
    ]

    for obj in objectives:
        p = objectives_tf.add_paragraph()
        p.text = obj
        p.font.size = Pt(12)

    # Components
    components_box = slide.shapes.add_textbox(Inches(5), Inches(1.3), Inches(4.25), Inches(2.5))
    components_tf = components_box.text_frame
    components_tf.text = "PROGRAM COMPONENTS"
    components_tf.paragraphs[0].font.bold = True
    components_tf.paragraphs[0].font.size = Pt(16)
    components_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    components = [
        "• NCD clinics at district level",
        "• Population-based screening",
        "• Free diagnostic services",
        "• Drug procurement & distribution",
        "• Health promotion activities"
    ]

    for comp in components:
        p = components_tf.add_paragraph()
        p.text = comp
        p.font.size = Pt(12)

    # Implementation levels
    levels_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
    levels_tf = levels_box.text_frame
    levels_tf.text = "IMPLEMENTATION: National → State → District → Community (ASHA/ANM workers)"
    levels_tf.paragraphs[0].font.size = Pt(14)
    levels_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    levels_tf.paragraphs[0].font.bold = True

def create_national_program_visual_slide(prs):
    """National program visual diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "NPCDCS Implementation Framework"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add national program diagram
    try:
        diagram_path = "TLM_Diabetes_Mellitus/visualizations/national_program_diagram.png"
        slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    except:
        # Fallback text
        fallback_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
        fallback_tf = fallback_box.text_frame
        fallback_tf.text = "National → State → District → Community\n\nPrimary Prevention: Health promotion\nSecondary Prevention: Early diagnosis\nTertiary Prevention: Complication management\n\nKey Components: NCD clinics, screening camps, training, drugs, IEC activities"
        for para in fallback_tf.paragraphs:
            para.font.size = Pt(16)

def create_diet_exercise_slide(prs):
    """Diet and exercise recommendations slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Diet & Physical Activity for Diabetes Control & Prevention"
    title_tf.paragraphs[0].font.size = Pt(28)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Dietary recommendations
    diet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(3))
    diet_tf = diet_box.text_frame
    diet_tf.text = "DIETARY RECOMMENDATIONS"
    diet_tf.paragraphs[0].font.bold = True
    diet_tf.paragraphs[0].font.size = Pt(18)
    diet_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    diet_points = [
        "• Mediterranean/DASH diets preferred",
        "• Low glycemic index carbohydrates",
        "• Increased fiber intake (25-30g/day)",
        "• Plant-based proteins (pulses, legumes)",
        "• Reduced saturated/trans fats (<7% energy)",
        "• Portion control and mindful eating"
    ]

    for point in diet_points:
        p = diet_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(12)

    # Exercise recommendations
    exercise_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(3))
    exercise_tf = exercise_box.text_frame
    exercise_tf.text = "PHYSICAL ACTIVITY PRESCRIPTION"
    exercise_tf.paragraphs[0].font.bold = True
    exercise_tf.paragraphs[0].font.size = Pt(18)
    exercise_tf.paragraphs[0].font.color.rgb = RGBColor(155, 89, 182)

    exercise_points = [
        "• 150 minutes/week moderate intensity",
        "• Combination of aerobic + resistance training",
        "• Yoga and flexibility exercises",
        "• Walking/bicycling for commuting",
        "• Sedentary behavior reduction",
        "• Individualized based on fitness level"
    ]

    for point in exercise_points:
        p = exercise_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(12)

    # Evidence statement
    evidence_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    evidence_tf = evidence_box.text_frame
    evidence_tf.text = "Evidence: DPP trial showed 58% reduction in diabetes incidence with lifestyle intervention alone."
    evidence_tf.paragraphs[0].font.size = Pt(14)
    evidence_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    evidence_tf.paragraphs[0].font.bold = True

def create_challenges_slide(prs):
    """Challenges and solutions slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Challenges in Diabetes Control & Prevention"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Left challenges
    challenges_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(3))
    challenges_tf = challenges_box.text_frame
    challenges_tf.text = "KEY CHALLENGES"
    challenges_tf.paragraphs[0].font.bold = True
    challenges_tf.paragraphs[0].font.size = Pt(18)
    challenges_tf.paragraphs[0].font.color.rgb = RGBColor(231, 76, 60)

    challenges = [
        "• Therapeutic inertia (delay in treatment intensification)",
        "• Adherence to lifestyle changes long-term",
        "• Hypoglycemia risk with intensive control",
        "• Resource constraints in low-income settings",
        "• Cultural and behavioral barriers",
        "• Socioeconomic disparities in access"
    ]

    for challenge in challenges:
        p = challenges_tf.add_paragraph()
        p.text = challenge
        p.font.size = Pt(12)

    # Right solutions
    solutions_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(3))
    solutions_tf = solutions_box.text_frame
    solutions_tf.text = "STRATEGIC SOLUTIONS"
    solutions_tf.paragraphs[0].font.bold = True
    solutions_tf.paragraphs[0].font.size = Pt(18)
    solutions_tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)

    solutions = [
        "• Regular treatment algorithm audits",
        "• Behavioral support and counseling",
        "• CGM-guided therapy adjustments",
        "• Task shifting to trained non-physicians",
        "• Culturally adapted health education",
        "• Community-based peer support groups"
    ]

    for solution in solutions:
        p = solutions_tf.add_paragraph()
        p.text = solution
        p.font.size = Pt(12)

def create_conclusion_slide(prs):
    """Conclusion and key takeaways slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Conclusion & Key Takeaways"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Key takeaways
    takeaways_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
    takeaways_tf = takeaways_box.text_frame
    takeaways_tf.text = "CONTROL STRATEGIES"
    takeaways_tf.paragraphs[0].font.bold = True
    takeaways_tf.paragraphs[0].font.size = Pt(20)
    takeaways_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    control_points = [
        "• Individualized glycemic targets based on patient characteristics (age, comorbidities, hypoglycemia risk)",
        "• Comprehensive monitoring: SMBG/CGM + HbA1c every 3-6 months",
        "• Lifestyle intervention as foundation + metformin first-line pharmacological therapy",
        "• Multidisciplinary team approach with regular complication screening",
        "• Technology utilization (CGM, pumps) and behavioral support for optimal outcomes"
    ]

    for point in control_points:
        p = takeaways_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        if "PREVENTION" in point or "NATIONAL" in point:
            p.font.bold = True

    prevention_title = takeaways_tf.add_paragraph()
    prevention_title.text = "PREVENTION FRAMEWORK"
    prevention_title.font.bold = True
    prevention_title.font.size = Pt(20)
    prevention_title.font.color.rgb = RGBColor(46, 204, 113)
    prevention_title.level = 0

    prevention_points = [
        "• Primary prevention: DPP model (58% reduction) - lifestyle intervention for high-risk individuals",
        "• Secondary prevention: Prediabetes management with lifestyle ± metformin",
        "• Tertiary prevention: Optimal control and complication management in established diabetes",
        "• Population-based approaches combined with high-risk individual strategies"
    ]

    for point in prevention_points:
        p = takeaways_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)

    national_title = takeaways_tf.add_paragraph()
    national_title.text = "NATIONAL PROGRAMME (NPCDCS)"
    national_title.font.bold = True
    national_title.font.size = Pt(20)
    national_title.font.color.rgb = RGBColor(155, 89, 182)
    national_title.level = 0

    national_points = [
        "• Comprehensive government initiative for diabetes and NCD control",
        "• Multi-level implementation: National → State → District → Community",
        "• Key components: Screening, treatment, health promotion, capacity building",
        "• Focus on prevention, early diagnosis, and comprehensive management"
    ]

    for point in national_points:
        p = takeaways_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)

def main():
    """Main function to create the PowerPoint presentation"""
    # Create presentation
    prs = Presentation()

    # Set slide size to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Create slides
    print("Creating title slide...")
    create_title_slide(prs)

    print("Creating overview slide...")
    create_overview_slide(prs)

    print("Creating control targets slide...")
    create_control_targets_slide(prs)

    print("Creating monitoring slide...")
    create_monitoring_slide(prs)

    print("Creating prevention levels slide...")
    create_prevention_levels_slide(prs)

    print("Creating control strategies visual slide...")
    create_control_strategies_visual_slide(prs)

    print("Creating national program overview slide...")
    create_national_program_slide(prs)

    print("Creating national program visual slide...")
    create_national_program_visual_slide(prs)

    print("Creating diet and exercise slide...")
    create_diet_exercise_slide(prs)

    print("Creating challenges slide...")
    create_challenges_slide(prs)

    print("Creating conclusion slide...")
    create_conclusion_slide(prs)

    # Save presentation
    output_path = "TLM_Diabetes_Mellitus/Diabetes_Control_Prevention_Presentation.pptx"
    prs.save(output_path)
    print(f"\n✅ PowerPoint presentation saved successfully!")
    print(f"📁 Location: {output_path}")
    print(f"📊 Slides created: {len(prs.slides)}")
    print(f"\n📋 Presentation covers:")
    print(f"• Diabetes Control Strategies (targets, monitoring, pyramid)")
    print(f"• Comprehensive Prevention Framework (primary/secondary/tertiary)")
    print(f"• India's National Programme (NPCDCS implementation)")
    print(f"• Implementation challenges and solutions")

if __name__ == "__main__":
    main()
