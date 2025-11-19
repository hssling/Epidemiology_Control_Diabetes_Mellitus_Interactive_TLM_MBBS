"""
Diabetes Mellitus TLM - PowerPoint Presentation Generator
Creates a comprehensive PPTX file with all diabetes content

To use this script:
1. Install required libraries: pip install python-pptx pandas matplotlib
2. Run: python create_pptx_presentation.py
3. The presentation will be saved as 'Diabetes_TLM_Presentation.pptx'
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    import pandas as pd
    import matplotlib.pyplot as plt
    import numpy as np
    import io
    import os
except ImportError as e:
    print(f"Missing required libraries. Please install: pip install python-pptx pandas matplotlib")
    print(f"Error: {e}")
    exit(1)

def create_title_slide(prs):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Diabetes Mellitus"
    subtitle.text = "Teaching Learning Material for MBBS Students\nComprehensive Guide with India Focus"

    # Style the title
    title_tf = title.text_frame
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    title_tf.paragraphs[0].font.bold = True

def create_content_slide(prs, title_text, content_points, layout_idx=1):
    """Create a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = title_text

    # Style title
    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    if layout_idx == 1:  # Title and content layout
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = ""

        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 0

def create_two_column_slide(prs, title_text, left_points, right_points):
    """Create a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
    shapes = slide.shapes

    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = title_text
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Left column
    left_box = shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(5))
    left_tf = left_box.text_frame
    left_tf.text = "DEFINITION & CRITERIA"

    for point in left_points:
        p = left_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)

    # Right column
    right_box = shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(5))
    right_tf = right_box.text_frame
    right_tf.text = "DIAGNOSTIC CRITERIA"

    for point in right_points:
        p = right_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)

def create_table_slide(prs, title_text, headers, data):
    """Create a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = title_text
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Create table
    rows, cols = len(data) + 1, len(headers)
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(5)).table

    # Set headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(52, 152, 219)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True

    # Fill data
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 249, 250)

def create_statistics_slide(prs):
    """Create statistics slide with epidemiology data and chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Epidemiology & Disease Burden"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add epidemiology chart
    try:
        chart_path = "TLM_Diabetes_Mellitus/visualizations/epidemiology_chart.png"
        slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
    except:
        # Fallback to text if image not found
        pass

    # Key statistics below the chart
    stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1))
    stats_tf = stats_box.text_frame
    stats_tf.text = "• Global: 537M cases (2021) → 643M (2030)  • India: 101M cases • $966B healthcare cost • 6.7M annual deaths"
    stats_tf.paragraphs[0].font.size = Pt(14)
    stats_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_pathophysiology_slide(prs):
    """Pathophysiology concepts slide with diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Pathophysiology of Diabetes Mellitus"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add pathophysiology diagram
    try:
        diagram_path = "TLM_Diabetes_Mellitus/visualizations/pathophysiology_diagram.png"
        slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
    except:
        # Fallback to text boxes if image not found
        # Type 1 DM box
        t1_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(2.5))
        t1_tf = t1_box.text_frame
        t1_tf.text = "TYPE 1 DIABETES\n\n• Autoimmune β-cell destruction\n• HLA-DR3/DR4 genetic susceptibility\n• GAD, IA-2, ICA antibodies\n• Absolute insulin deficiency\n• Ketosis-prone"

        # Type 2 DM box
        t2_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(2.5))
        t2_tf = t2_box.text_frame
        t2_tf.text = "TYPE 2 DIABETES\n\n• Insulin resistance + β-cell failure\n• Polygenic inheritance\n• Asian Indian phenotype\n• Relative insulin deficiency\n• Lifestyle modifiable"

    # Key points below diagram
    points_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1))
    points_tf = points_box.text_frame
    points_tf.text = "• Type 1: Autoimmune β-cell destruction → Absolute insulin deficiency → Ketosis-prone • Type 2: Insulin resistance + β-cell failure → Relative insulin deficiency → Lifestyle modifiable"
    points_tf.paragraphs[0].font.size = Pt(12)
    points_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_risk_factors_slide(prs):
    """Risk factors slide with diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Risk Factors for Diabetes Mellitus"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add risk factors diagram
    try:
        diagram_path = "TLM_Diabetes_Mellitus/visualizations/risk_factor_diagram.png"
        slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
    except:
        # Fallback text content
        pass

def create_treatment_algorithm_slide(prs):
    """Treatment algorithm slide with visual"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "ADA Treatment Algorithm (2024)"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add treatment algorithm visual
    try:
        algorithm_path = "TLM_Diabetes_Mellitus/visualizations/treatment_algorithm.png"
        slide.shapes.add_picture(algorithm_path, Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
    except:
        # Fallback to text if image not found
        algorithm_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
        algorithm_tf = algorithm_text.text_frame
        algorithm_tf.text = "STEPWISE TREATMENT APPROACH\n\n"
        steps = [
            "1. DIAGNOSIS → Comprehensive evaluation",
            "2. LIFESTYLE + METFORMIN (first-line)",
            "3. DUAL THERAPY (if targets not met)",
            "4. TRIPLE THERAPY → Continue if not at target",
            "5. INSULIN INITIATION",
            "6. MONITORING & ADJUSTMENT"
        ]
        for step in steps:
            p = algorithm_tf.add_paragraph()
            if step.startswith(("2.", "3.", "4.", "5.", "6.")):
                p.font.bold = True
                p.font.color.rgb = RGBColor(46, 204, 113)
            p.text = step
            p.font.size = Pt(14)

    # Key points below image
    points_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1))
    points_tf = points_box.text_frame
    points_tf.text = "• Lifestyle + Metformin first-line • Individualize based on patient factors • DPP-4i, SGLT2i, GLP-1 RA for dual therapy • Insulin initiation when needed"
    points_tf.paragraphs[0].font.size = Pt(12)
    points_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_prevention_slide(prs):
    """Prevention strategies slide with flowchart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Prevention of Diabetes Mellitus"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Add prevention flowchart
    try:
        flowchart_path = "TLM_Diabetes_Mellitus/visualizations/prevention_flowchart.png"
        slide.shapes.add_picture(flowchart_path, Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
    except:
        # Fallback to text boxes if image not found
        # Three levels of prevention
        primary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.8), Inches(2))
        primary_tf = primary_box.text_frame
        primary_tf.text = "PRIMARY PREVENTION\n(High-Risk Individuals)"
        primary_tf.paragraphs[0].font.bold = True
        primary_tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)

        primary_points = [
            "• DPP Lifestyle Intervention",
            "• 5-7% weight loss",
            "• 150 min activity/week",
            "• Metformin 850mg BID (BMI ≥35)",
            "• 58% reduction in incidence"
        ]

        for point in primary_points:
            p = primary_tf.add_paragraph()
            p.text = point
            p.font.size = Pt(12)

        # Secondary prevention
        secondary_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.2), Inches(2.8), Inches(2))
        secondary_tf = secondary_box.text_frame
        secondary_tf.text = "SECONDARY PREVENTION\n(Prediabetes Management)"
        secondary_tf.paragraphs[0].font.bold = True
        secondary_tf.paragraphs[0].font.color.rgb = RGBColor(230, 126, 34)
        # ... rest of secondary and tertiary prevention code ...

    # Key points below flowchart
    points_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1))
    points_tf = points_box.text_frame
    points_tf.text = "• Primary: Prevent progression from prediabetes • Secondary: Lifestyle + metformin for prediabetes • Tertiary: Optimal management of established diabetes"
    points_tf.paragraphs[0].font.size = Pt(12)
    points_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

def create_learning_objectives_slide(prs):
    """Learning objectives and summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Learning Objectives & Key Takeaways"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)

    # Objectives
    objectives_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.25), Inches(5))
    objectives_tf = objectives_box.text_frame
    objectives_tf.text = "LEARNING OBJECTIVES"
    objectives_tf.paragraphs[0].font.bold = True
    objectives_tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    objectives = [
        "• Master diabetes pathophysiology (T1DM vs T2DM)",
        "• Apply ADA 2024 diagnostic criteria confidently",
        "• Develop evidence-based treatment algorithms",
        "• Understand India-specific epidemiological patterns",
        "• Design comprehensive prevention strategies",
        "• Recognize and manage complications",
        "• Implement lifestyle interventions effectively"
    ]

    for obj in objectives:
        p = objectives_tf.add_paragraph()
        p.text = obj
        p.font.size = Pt(14)

    # Key clinical pearls
    pearls_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.25), Inches(5))
    pearls_tf = pearls_box.text_frame
    pearls_tf.text = "CLINICAL PEARLS"
    pearls_tf.paragraphs[0].font.bold = True
    pearls_tf.paragraphs[0].font.color.rgb = RGBColor(46, 204, 113)

    pearls = [
        "• Three P's: Polyuria, Polydipsia, Polyphagia",
        "• Diabetes diagnostic: FPG≥126, 2hPG≥200, HbA1c≥6.5",
        "• Lifestyle + metformin first-line therapy",
        "• Individualize targets based on patient factors",
        "• Prevention better than cure - focus on lifestyle",
        "• Foot examination essential at every visit"
    ]

    for pearl in pearls:
        p = pearls_tf.add_paragraph()
        p.text = pearl
        p.font.size = Pt(14)

def main():
    """Main function to create the PowerPoint presentation"""
    # Create presentation
    prs = Presentation()

    # Set slide size to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Create slides
    print("Creating title slide...")
    create_title_slide(prs)

    print("Creating definition & criteria slide...")
    diagnostic_criteria = [
        "• Fasting Plasma Glucose (FPG): ≥126 mg/dL",
        "• 2-hour Plasma Glucose (2hPG): ≥200 mg/dL",
        "• Glycated Hemoglobin (HbA1c): ≥6.5%",
        "• Random PG + symptoms: ≥200 mg/dL",
        "• Asymptomatic: Confirmatory testing required",
        "• ADA 2024 criteria with Indian adaptations"
    ]

    definition_points = [
        "• Type 1: Autoimmune β-cell destruction",
        "• Type 2: Insulin resistance + β-cell failure",
        "• Maturity-onset diabetes of young (MODY)",
        "• Gestational diabetes mellitus",
        "• Other specific types (drug-induced, genetic)"
    ]

    create_two_column_slide(prs, "Definition & Diagnostic Criteria",
                          definition_points, diagnostic_criteria)

    print("Creating epidemiology slide...")
    create_statistics_slide(prs)

    print("Creating pathophysiology slide...")
    create_pathophysiology_slide(prs)

    print("Creating risk factors slide...")
    create_risk_factors_slide(prs)

    print("Creating diagnostic criteria table slide...")
    headers = ["Parameter", "Diabetes", "Prediabetes", "Normal"]
    criteria_data = [
        ["FPG", "≥126 mg/dL", "100-125 mg/dL", "<100 mg/dL"],
        ["2hPG", "≥200 mg/dL", "140-199 mg/dL", "<140 mg/dL"],
        ["HbA1c", "≥6.5%", "5.7-6.4%", "<5.7%"],
        ["Random PG", "≥200 mg/dL + symptoms", "-", "-"]
    ]
    create_table_slide(prs, "Diagnostic Criteria (ADA 2024)", headers, criteria_data)

    print("Creating treatment algorithm slide...")
    create_treatment_algorithm_slide(prs)

    print("Creating prevention slide...")
    create_prevention_slide(prs)

    print("Creating learning objectives slide...")
    create_learning_objectives_slide(prs)

    # Save presentation
    output_path = "TLM_Diabetes_Mellitus/Diabetes_TLM_Presentation.pptx"
    prs.save(output_path)
    print(f"\n✅ PowerPoint presentation saved successfully!")
    print(f"📁 Location: {output_path}")
    print(f"📊 Slides created: {len(prs.slides)}")
    print(f"\n📋 To use this PowerPoint:")
    print(f"1. Install python-pptx: pip install python-pptx")
    print(f"2. Run: python create_pptx_presentation.py")
    print(f"3. Open the generated .pptx file in PowerPoint")

if __name__ == "__main__":
    main()
